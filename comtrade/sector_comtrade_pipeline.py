"""
sector_comtrade_pipeline.py
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Reusable config-driven Comtrade trade-data pipeline.
All parameters (sector, HS codes, years, flows, buckets) come
from a sector config Excel file — no manual code editing needed.

USAGE
─────
  Setup check  : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode setup_check
  Test run     : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode test
  Full pull    : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode full
  Clean        : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode clean
  Slide-ready      : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode slide_ready
  Competitor test  : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode competitor_test
  Competitor full  : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode competitor_full
  Visuals          : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode visuals
  All modes        : python sector_comtrade_pipeline.py --config textiles_sector_config.xlsx --mode all

API KEY (set before running)
─────────────────────────────
  Mac/Linux : export COMTRADE_API_KEY=your_key_here
  Windows   : $env:COMTRADE_API_KEY="your_key_here"

DO NOT OVERWRITE comtrade.py or any original organisational files.
"""

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# IMPORTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
import argparse
import logging
import re
import os
import sys
import time
import traceback
from datetime import datetime
from pathlib import Path

import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from tqdm import tqdm

import comtradeapicall

# In-memory cache for the session
_COUNTRY_MAP = {}

def get_country_mapping():
    """Fetches and caches the official Comtrade reporter mapping."""
    global _COUNTRY_MAP
    if not _COUNTRY_MAP:
        try:
            # Returns a list of dicts: [{'id': 842, 'text': 'USA'}, ...]
            refs = comtradeapicall.getReference('reporter')
            # Create a mapping of "Country Name" -> integer ID
            _COUNTRY_MAP = {str(item['text']): int(item['id']) for item in refs if item['id'] != 'all'}
        except Exception as e:
            log.error(f"Failed to fetch country references: {e}")
            return {}
    return _COUNTRY_MAP

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# AUTO-LOAD .env FILE
# Looks for .env in the same folder as this script.
# Keys already set in the environment are not overwritten.
# The key value is never printed.
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _load_dotenv():
    env_path = Path(__file__).parent / ".env"
    if not env_path.exists():
        return
    with open(env_path, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            key, _, val = line.partition("=")
            key = key.strip()
            val = val.strip().strip('"').strip("'")
            if key and key not in os.environ:   # don't overwrite if already set
                os.environ[key] = val

_load_dotenv()

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# LOGGING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
_errors_buffer = []   # runtime errors collected during a run

def setup_logging(log_file="pipeline.log"):
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s  %(levelname)-8s  %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
        handlers=[
            logging.FileHandler(log_file, encoding="utf-8"),
            logging.StreamHandler(sys.stdout),
        ],
    )

log = logging.getLogger(__name__)


def record_error(mode, flow, hs_code, year, etype, emsg, action):
    _errors_buffer.append({
        "Timestamp":     datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "Mode":          mode,
        "Flow":          flow,
        "HS Code":       hs_code,
        "Year":          year,
        "Error type":    etype,
        "Error message": emsg,
        "Action taken":  action,
    })


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CONFIG READER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def read_config(config_file):
    """Parse the sector config Excel file into a structured dict."""
    xl = pd.ExcelFile(config_file)
    # All config sheets have: row 1 = title, row 2 = header, row 3+ = data → header=1 (0-indexed)

    # sector_details
    df_d = pd.read_excel(xl, "sector_details", header=1)
    df_d.columns = [str(c).strip() for c in df_d.columns]
    sector = dict(zip(df_d.iloc[:, 0].astype(str).str.strip(),
                      df_d.iloc[:, 1].astype(str).str.strip()))

    # hs_codes  — dtype=str preserves leading zeros
    df_hs = pd.read_excel(xl, "hs_codes", header=1, dtype=str)
    df_hs.columns = [str(c).strip() for c in df_hs.columns]
    df_hs = df_hs.dropna(subset=["HS Code"])
    df_hs["HS Code"] = df_hs["HS Code"].str.strip().str.zfill(2)
    all_hs = df_hs
    # Accept either column name (old template vs new QE-generated config)
    incl_col = next((c for c in df_hs.columns if "include" in c.lower()), None)
    if incl_col is None:
        raise KeyError("Could not find Include column in hs_codes sheet. Expected 'Include/Review/Exclude' or 'Include / Exclude'.")
    included_hs = df_hs[df_hs[incl_col].str.strip().str.lower() == "include"]["HS Code"].tolist()

    # sector_buckets  — build code -> bucket map
    df_b = pd.read_excel(xl, "sector_buckets", header=1)
    df_b.columns = [str(c).strip() for c in df_b.columns]
    df_b = df_b.dropna(subset=["Sector Bucket"])
    # Drop the trailing instructional/placeholder row every generated template
    # includes (e.g. "← Fill in the Definition column..."). It has text in
    # "Sector Bucket" so the dropna above doesn't catch it.
    df_b = df_b[~df_b["Sector Bucket"].astype(str).str.strip().str.startswith("←")]
    # Accept either column name for HS codes in bucket sheet
    hs_col = next((c for c in df_b.columns if "hs" in c.lower() and ("code" in c.lower() or "incl" in c.lower())), None)
    if hs_col is None:
        raise KeyError("Could not find HS codes column in sector_buckets sheet. Expected 'Included HS codes' or 'HS Codes'.")
    bucket_map = {}
    for _, row in df_b.iterrows():
        bname = str(row["Sector Bucket"]).strip()
        codes_str = str(row[hs_col]).strip()
        if not codes_str or codes_str.lower() in ("na", "nan", "none"):
            continue
        for code in codes_str.split(","):
            code = code.strip()
            if not code or code.lower() in ("na", "nan", "none"):
                continue
            # Normalize to a 2-digit HS chapter regardless of source code length
            # (sector_buckets may have been populated at HS level 2, 4 or 6 —
            # build_master_table always looks Sector Bucket up by the 2-digit
            # chapter, so the map must always be keyed by 2-digit chapters too;
            # zfill alone does not truncate longer codes, which previously
            # caused every row to silently fall back to "Other").
            code = code.zfill(2)[:2]
            bucket_map[code] = bname

    # years
    import datetime as _dt
    df_y = pd.read_excel(xl, "years", header=1)
    df_y.columns = [str(c).strip() for c in df_y.columns]
    df_y = df_y.dropna(subset=["Start Year"])
    start_year   = int(float(df_y.iloc[0]["Start Year"]))

    # End Year: use config value if set, otherwise auto-detect as (current year - 1)
    # Comtrade annual data typically lags ~12 months, so current_year-1 is the
    # latest year that is reliably available.
    _ey_raw = df_y.iloc[0].get("End Year", None)
    if pd.isna(_ey_raw) or str(_ey_raw).strip().lower() in ("", "auto", "latest", "nan"):
        end_year = _dt.date.today().year - 1
        log.info(f"  ·  End Year auto-detected: {end_year} (current year − 1)")
    else:
        end_year = int(float(_ey_raw))

    # Always rebuild period string from start→end (ignore any hardcoded value)
    years_list = list(range(start_year, end_year + 1))
    period     = ",".join(str(y) for y in years_list)

    # flows
    df_f = pd.read_excel(xl, "flows", header=1)
    df_f.columns = [str(c).strip() for c in df_f.columns]
    df_f = df_f.dropna(subset=["Flow Name"])
    # Drop the trailing instructional/placeholder row that every generated config
    # template includes (e.g. "← Usually no changes needed..."). That row has text
    # in "Flow Name" so the dropna above doesn't catch it, but it has no real Flow
    # Code — a genuine flow always specifies "X" (exports) or "M" (imports). Without
    # this filter the pipeline previously tried to fetch a 5th, bogus flow on every
    # single run (output filename "TEST_nan" / "nan"), which always failed and
    # polluted the error log and Errors_Log sheet for every sector.
    df_f = df_f[df_f["Flow Code"].notna()]
    df_f = df_f[~df_f["Flow Name"].astype(str).str.strip().str.startswith("←")]
    # Accept "Output File Name" or "Output File"
    out_col = next((c for c in df_f.columns if c.lower().startswith("output")), "Output File Name")
    flows = []
    for _, row in df_f.iterrows():
        rc = str(row["Reporter Code"]).strip()
        reporter_code = None if rc.lower() in ("none", "all", "nan", "") else int(float(rc))
        pc = str(row["Partner Code"]).strip()
        partner_code = str(int(float(pc))) if pc not in ("nan", "") else "0"
        fname = str(row[out_col]).strip()
        flows.append({
            "label":        str(row["Flow Name"]).strip(),
            "flowCode":     str(row["Flow Code"]).strip(),
            "reporterCode": reporter_code,
            "partnerCode":  partner_code,
            "output":       fname,
            "test_output":  "TEST_" + fname,
        })

    # settings
    df_s = pd.read_excel(xl, "settings", header=1)
    df_s.columns = [str(c).strip() for c in df_s.columns]
    df_s = df_s.dropna(subset=["Setting"])
    settings = dict(zip(df_s.iloc[:, 0].astype(str).str.strip(),
                        df_s.iloc[:, 1].astype(str).str.strip()))

    slug               = sector.get("Sector slug", "unknown").strip()
    api_key_var        = settings.get("API key environment variable name", "COMTRADE_API_KEY").strip()
    test_hs            = str(settings.get("Test HS code", included_hs[0] if included_hs else "61")).strip().zfill(2)
    test_year          = str(settings.get("Test year", str(end_year))).strip()
    output_folder      = settings.get("Output folder", ".").strip()
    error_log_file     = settings.get("Error log filename", f"{slug}_errors.log").strip()
    continue_on_error  = settings.get("Continue on error", "Yes").strip().lower() == "yes"
    overwrite_existing = settings.get("Overwrite existing output files", "Yes").strip().lower() == "yes"

    # ── optional: dgcis_mapping sheet ───────────────────────
    dgcis_config = None
    if "dgcis_mapping" in xl.sheet_names:
        df_dm = pd.read_excel(xl, "dgcis_mapping", header=1)
        df_dm.columns = [str(c).strip() for c in df_dm.columns]
        df_dm = df_dm.dropna(how="all")
        dgcis_config = {}
        for _, row in df_dm.iterrows():
            key = str(row.iloc[0]).strip()
            val = str(row.iloc[1]).strip() if len(row) > 1 else ""
            dgcis_config[key] = val

    # ── optional: competitor_settings sheet ───────────────
    competitor_settings = {
        "Top exporters to retain": "10",
        "Top importers to retain": "10",
        "Analyse by overall sector": "Yes",
        "Analyse by sector bucket": "Yes",
        "Analyse by HS6": "No",
        "Minimum trade value threshold (USD Mn)": "0",
        "Years for market share trend": "",
    }
    if "competitor_settings" in xl.sheet_names:
        df_cs = pd.read_excel(xl, "competitor_settings", header=1)
        df_cs.columns = [str(c).strip() for c in df_cs.columns]
        df_cs = df_cs.dropna(how="all")
        for _, row in df_cs.iterrows():
            k = str(row.iloc[0]).strip()
            v = str(row.iloc[1]).strip() if len(row) > 1 else ""
            if k and k != "nan":
                competitor_settings[k] = v

    # ── optional: visual_settings sheet ───────────────────
    visual_settings = {
        "Reference PowerPoint file name": "",
        "Output PowerPoint file name": f"{slug}_context_setting_graphs.pptx",
        "Output chart data workbook name": f"{slug}_context_graph_data.xlsx",
        "Create placeholder slides for missing charts": "Yes",
        "Source note format": f"UN Comtrade API, {{years}}. All values USD Mn.",
    }
    if "visual_settings" in xl.sheet_names:
        df_vs = pd.read_excel(xl, "visual_settings", header=1)
        df_vs.columns = [str(c).strip() for c in df_vs.columns]
        df_vs = df_vs.dropna(how="all")
        for _, row in df_vs.iterrows():
            k = str(row.iloc[0]).strip()
            v = str(row.iloc[1]).strip() if len(row) > 1 else ""
            if k and k != "nan":
                visual_settings[k] = v

    return {
        "sector":               sector,
        "slug":                 slug,
        "sector_name":          sector.get("Sector name", slug),
        "all_hs_df":            all_hs,
        "hs_codes":             included_hs,
        "bucket_map":           bucket_map,
        "bucket_df":            df_b.rename(columns={hs_col: "Included HS codes"}) if hs_col != "Included HS codes" else df_b,
        "period":               period,
        "start_year":           start_year,
        "end_year":             end_year,
        "years_list":           years_list,
        "flows":                flows,
        "settings":             settings,
        "api_key_var":          api_key_var,
        "test_hs":              test_hs,
        "test_year":            test_year,
        "output_folder":        output_folder,
        "error_log_file":       error_log_file,
        "continue_on_error":    continue_on_error,
        "overwrite":            overwrite_existing,
        "cleaned_file":         f"{slug}_comtrade_cleaned.xlsx",
        "slide_ready_file":     f"{slug}_slide_ready_tables.xlsx",
        "pptx_file":            f"{slug}_trade_landscape_draft.pptx",
        "dgcis_config":         dgcis_config,
        "competitor_settings":  competitor_settings,
        "visual_settings":      visual_settings,
    }


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SETUP CHECK
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def run_setup_check(config):
    """Check Python version, libraries, config validity and API key."""
    log.info("=" * 60)
    log.info("  SETUP CHECK")
    log.info("=" * 60)

    checks = []

    # Python version
    pv = sys.version_info
    checks.append(("Python version", f"{pv.major}.{pv.minor}.{pv.micro}",
                    pv >= (3, 8), "Required: 3.8+"))

    # Required libraries
    for lib in ["comtradeapicall", "pandas", "openpyxl", "tqdm"]:
        try:
            __import__(lib)
            checks.append((f"Library: {lib}", "Installed", True, ""))
        except ImportError:
            checks.append((f"Library: {lib}", "NOT INSTALLED", False, f"Run: pip install {lib}"))

    # Optional PPTX library
    try:
        import pptx  # noqa
        checks.append(("Library: python-pptx (optional)", "Installed", True, ""))
    except ImportError:
        checks.append(("Library: python-pptx (optional)", "Not installed",
                        True, "Install for PPTX output: pip install python-pptx"))

    # API key
    api_key = os.environ.get(config["api_key_var"], "")
    has_key = bool(api_key) and api_key != "YOUR_API_KEY_HERE"
    checks.append(("API key env var",
                    f'{config["api_key_var"]} {"set ✓" if has_key else "NOT SET ✗"}',
                    has_key,
                    f"Export: export {config['api_key_var']}=your_key_here"))

    # Config validity
    checks.append(("Sector name",    config["sector_name"],                  bool(config["sector_name"]), ""))
    checks.append(("Sector slug",    config["slug"],                          bool(config["slug"]), ""))
    n = len(config["hs_codes"])
    checks.append(("HS codes (Include)", str(n), n > 0, "Mark at least one row as Include"))
    checks.append(("HS codes are strings", "Yes" if all(isinstance(c, str) for c in config["hs_codes"]) else "No",
                    all(isinstance(c, str) for c in config["hs_codes"]), "Check hs_codes sheet"))
    checks.append(("Leading zeros preserved",
                    str(config["hs_codes"][:5]),
                    True, ""))
    checks.append(("Period string",  config["period"], bool(config["period"]), ""))
    checks.append(("Number of flows", str(len(config["flows"])), len(config["flows"]) > 0, ""))
    checks.append(("Number of buckets", str(len(config["bucket_map"])), len(config["bucket_map"]) > 0, ""))
    checks.append(("Test HS code",   config["test_hs"], bool(config["test_hs"]), ""))
    checks.append(("Output folder",  config["output_folder"], True, ""))

    all_pass = True
    for item, value, passed, note in checks:
        status = "✓" if passed else "✗"
        line = f"  {status}  {item:<45} {value}"
        if not passed:
            line += f"  ← {note}"
            all_pass = False
        log.info(line)

    if all_pass:
        log.info("\n  ✅  All checks passed. Ready to run.\n")
    else:
        log.info("\n  ⚠   Some checks failed. Fix the issues above before running.\n")

    return all_pass


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# API FETCH
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def fetch_flow(config, flow_cfg, hs_codes, period, is_test, mode_label):
    """Fetch Comtrade data for one flow. Returns (DataFrame, failed_codes)."""
    import comtradeapicall

    api_key     = os.environ.get(config["api_key_var"], "")
    label       = flow_cfg["label"]
    flow_code   = flow_cfg["flowCode"]
    reporter    = flow_cfg["reporterCode"]
    partner     = flow_cfg["partnerCode"]
    output_file = flow_cfg["test_output"] if is_test else flow_cfg["output"]

    log.info("=" * 60)
    log.info(f"  FLOW: {label}{'  [TEST]' if is_test else ''}")
    log.info(f"  Period:   {period}")
    log.info(f"  Reporter: {reporter if reporter is not None else 'All (None)'}")
    log.info(f"  Partner:  {partner}  |  FlowCode: {flow_code}")
    log.info(f"  HS codes: {hs_codes}")
    log.info(f"  Output:   {output_file}")
    log.info("=" * 60)

    results, failed = [], []

    for code in tqdm(hs_codes, desc=label):
        log.info(f"  → HS {code} ...")
        try:
            mydf = comtradeapicall.getFinalData(
                api_key,
                typeCode      = "C",
                freqCode      = "A",
                clCode        = "HS",
                period        = period,
                reporterCode  = reporter,
                cmdCode       = code,
                flowCode      = flow_code,
                partnerCode   = partner,
                partner2Code  = None,
                customsCode   = None,
                motCode       = None,
                maxRecords    = None,
                format_output = "JSON",
                aggregateBy   = None,
                breakdownMode = "classic",
                countOnly     = None,
                includeDesc   = True,
            )
            if mydf is None or mydf.empty:
                log.warning(f"     ⚠  No data returned for HS {code}")
                failed.append(code)
                record_error(mode_label, label, code, period, "NoData",
                             "API returned empty/None", "Skipped — logged")
            else:
                log.info(f"     ✓  {len(mydf)} rows for HS {code}")
                results.append(mydf)
        except Exception:
            emsg = traceback.format_exc().strip().split("\n")[-1]
            log.error(f"     ✗  ERROR for HS {code}: {emsg}")
            failed.append(code)
            record_error(mode_label, label, code, period, "APIError", emsg,
                         "Skipped — logged" if config["continue_on_error"] else "Halted")
            if not config["continue_on_error"]:
                log.error("Continue-on-error is OFF. Stopping flow.")
                break
        time.sleep(0.5)

    combined = pd.concat(results, ignore_index=True, sort=False) if results else pd.DataFrame()
    if not combined.empty:
        combined.to_excel(output_file, index=False)
        log.info(f"  ✅  Saved {len(combined)} rows → {output_file}")
    else:
        log.warning(f"  ⚠  No data for {label}. Output file not created.")

    if failed:
        log.warning(f"  HS codes with no data/errors: {failed}")
    return combined, failed

def fetch_bilateral(config, reporter_name, partner_name, flow_code, hs_codes, period):
    api_key = os.environ.get(config["api_key_var"], "")
    cmap = get_country_mapping()
    
    reporter_code = cmap.get(reporter_name)
    partner_code = cmap.get(partner_name)
    
    if not reporter_code or not partner_code:
        log.error("Invalid country names provided for bilateral fetch.")
        return pd.DataFrame()

    log.info(f"Bilateral Fetch: {reporter_name} to {partner_name} | Flow: {flow_code} | Period: {period}")
    
    results = []
    for code in tqdm(hs_codes, desc=f"Bilateral {reporter_name}-{partner_name}"):
        try:
            # Reusing getFinalData, but with strict partner codes
            mydf = comtradeapicall.getFinalData(
                api_key, typeCode="C", freqCode="A", clCode="HS", period=period,
                reporterCode=reporter_code, cmdCode=code, flowCode=flow_code,
                partnerCode=partner_code, partner2Code=None, customsCode=None,
                motCode=None, maxRecords=None, format_output="JSON", aggregateBy=None,
                breakdownMode="classic", countOnly=None, includeDesc=True
            )
            if mydf is not None and not mydf.empty:
                results.append(mydf)
        except Exception as e:
            emsg = traceback.format_exc().strip().split("\n")[-1]
            record_error("bilateral", flow_code, code, period, "APIError", emsg, "Skipped")
            
        time.sleep(0.5) # Respect rate limits
        
    if not results:
        return pd.DataFrame()

    df = pd.concat(results, ignore_index=True)
    
    # 1. Clean and Calculate Unit Price (Reusing your logic)
    trade_usd = pd.to_numeric(df.get("primaryValue"), errors="coerce")
    qty_raw = pd.to_numeric(df.get("qty"), errors="coerce")
    qty_valid = qty_raw.where(qty_raw > 0)
    
    clean_df = pd.DataFrame({
        "Year": pd.to_numeric(df.get("period"), errors="coerce"),
        "Reporter": reporter_name,
        "Partner": partner_name,
        "Flow": "Export" if flow_code == "X" else "Import",
        "HS Code": df.get("cmdCode").astype(str).str.zfill(2),
        "HS Description": df.get("cmdDesc"),
        "Trade Value (USD)": trade_usd,
        "Trade Value USD Mn": (trade_usd / 1_000_000).round(4),
        "Quantity": qty_valid,
        "Qty Unit": df.get("qtyUnitAbbr").astype(str).where(qty_valid.notna(), ""),
        "Unit Price (USD)": (trade_usd / qty_valid).round(4)
    })
    
    # 2. Append to Geography_Bilateral Sheet
    _append_to_bilateral_sheet(config["cleaned_file"], clean_df)
    
    return clean_df

def _append_to_bilateral_sheet(filepath, new_data):
    """Appends to Geography_Bilateral, deduplicating existing records."""
    sheet_name = "Geography_Bilateral"
    if not os.path.exists(filepath):
        return

    try:
        # Load existing data if sheet exists
        existing_df = pd.DataFrame()
        xls = pd.ExcelFile(filepath)
        if sheet_name in xls.sheet_names:
            existing_df = pd.read_excel(xls, sheet_name=sheet_name, dtype={"HS Code": str})
        
        # Combine and deduplicate
        combined = pd.concat([existing_df, new_data], ignore_index=True)
        combined.drop_duplicates(
            subset=["Reporter", "Partner", "Flow", "HS Code", "Year"], 
            keep="last", inplace=True
        )
        
        # Save back to workbook
        with pd.ExcelWriter(filepath, engine="openpyxl", mode="a", if_sheet_exists="replace") as writer:
            combined.to_excel(writer, sheet_name=sheet_name, index=False)
            
    except Exception as e:
        log.error(f"Failed to update Geography_Bilateral sheet: {e}")


def run_fetch(config, is_test):
    """Run all configured flows and print a summary."""
    mode_label = "test" if is_test else "full"
    hs_codes   = [config["test_hs"]] if is_test else config["hs_codes"]
    period     = config["test_year"] if is_test else config["period"]

    header = "TEST MODE — 1 HS code, 1 year" if is_test else \
             f"FULL MODE — {len(hs_codes)} HS codes, years {config['start_year']}–{config['end_year']}"
    log.info("\n" + "█" * 60)
    log.info(f"  {header}")
    log.info("█" * 60)

    summary = []
    for flow in config["flows"]:
        df, failed = fetch_flow(config, flow, hs_codes, period, is_test, mode_label)
        summary.append({
            "Flow":         flow["label"],
            "Rows fetched": len(df),
            "Failed codes": failed or "None",
            "Output":       flow["test_output"] if is_test else flow["output"],
        })

    log.info("\n" + "=" * 60 + "\n  RUN SUMMARY\n" + "=" * 60)
    for s in summary:
        log.info(f"  {s['Flow']:<22} | {s['Rows fetched']:>7} rows | failed: {str(s['Failed codes']):<15} | → {s['Output']}")
    log.info("=" * 60)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CLEAN MODE — DATA HELPERS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _get(df, col_variants, fallback=""):
    """Case-insensitive column lookup."""
    cols_lower = {c.lower(): c for c in df.columns}
    for v in col_variants:
        if v.lower() in cols_lower:
            return df[cols_lower[v.lower()]]
    return pd.Series([fallback] * len(df))


def load_output_files(config):
    """Load the 4 full output xlsx files. Returns {label: DataFrame}."""
    frames = {}
    for flow in config["flows"]:
        path = flow["output"]
        if not os.path.exists(path):
            log.warning(f"  ⚠  Output file not found, skipping: {path}")
            record_error("clean", flow["label"], "ALL", "ALL", "FileNotFound",
                         f"{path} does not exist", "Skipped")
            continue
        try:
            df = pd.read_excel(path, dtype={"cmdCode": str})
            df["_source_file"]  = path
            df["_flow_label"]   = flow["label"]
            frames[flow["label"]] = df
            log.info(f"  ✓  Loaded {len(df)} rows from {path}")
        except Exception:
            emsg = traceback.format_exc().strip().split("\n")[-1]
            log.error(f"  ✗  Could not read {path}: {emsg}")
            record_error("clean", flow["label"], "ALL", "ALL", "ReadError", emsg, "Skipped")
    return frames


def build_master_table(frames, config):
    """Combine all frames into the Master_Table with standardised columns."""
    bucket_map = config["bucket_map"]
    all_rows   = []

    for label, df in frames.items():
        cmd_raw  = _get(df, ["cmdCode", "cmdcode"]).astype(str).str.zfill(2)
        chapter  = cmd_raw.str[:2]
        trade_usd = pd.to_numeric(_get(df, ["primaryValue", "primaryvalue"]), errors="coerce")

        # Quantity / unit price. Comtrade reports qty=0 (not NaN/blank) for the
        # large share of records that only carry a trade value with no
        # associated unit-quantity (weight-only or value-only records) — see
        # altQty/netWgt for those. Dividing by that zero would produce inf, so
        # qty<=0 is treated the same as missing: Quantity/Unit Price stay blank
        # rather than showing a bogus $0.00-per-unit or infinite price.
        qty_raw   = pd.to_numeric(_get(df, ["qty"]), errors="coerce")
        qty_valid = qty_raw.where(qty_raw > 0)
        qty_unit  = _get(df, ["qtyUnitAbbr", "qtyunitabbr"]).astype(str).where(qty_valid.notna(), "")
        unit_price = (trade_usd / qty_valid).round(4)

        row_df = pd.DataFrame({
            "HS Code":             cmd_raw,
            # Cast to str defensively, same as "HS Code" above: if a source
            # workbook's cmdDesc column ends up all-numeric (e.g. re-saved
            # by Excel, or a description that happens to look like a bare
            # number), pandas infers a numeric dtype on read, and every
            # downstream `.str[:55]` call on "Product Description" (6 call
            # sites) would crash with AttributeError instead of degrading
            # gracefully.
            "Product Description": _get(df, ["cmdDesc", "cmddesc"]).astype(str),
            "Sector Bucket":       chapter.map(bucket_map).fillna("Other"),
            "Reporter":            _get(df, ["reporterDesc", "reporterdesc"]),
            "Partner":             _get(df, ["partnerDesc", "partnerdesc"]),
            "Flow":                _get(df, ["flowDesc", "flowdesc"]).fillna(label),
            "Year":                pd.to_numeric(_get(df, ["period"]), errors="coerce"),
            "Trade Value USD":     trade_usd,
            "Trade Value USD Mn":  (trade_usd / 1_000_000).round(4),
            "Quantity":            qty_valid,
            "Qty Unit":            qty_unit,
            "Unit Price (USD)":    unit_price,
            "Source File":         df["_source_file"],
            "Source Database":     "UN Comtrade",
            "Notes":               "",
        })
        all_rows.append(row_df)

    if not all_rows:
        return pd.DataFrame()
    return pd.concat(all_rows, ignore_index=True)


def build_pivot(master, source_key, config):
    """Pivot one flow by (HS Code, Product Description, Sector Bucket) × Year."""
    subset = master[master["Source File"].str.contains(source_key, case=False, na=False)].copy()
    if subset.empty:
        return pd.DataFrame()

    subset["Chapter"] = subset["HS Code"].str[:2]
    grp = (subset.groupby(["Chapter", "Sector Bucket", "Year"], as_index=False)["Trade Value USD Mn"]
           .sum().round(2))

    desc_map = (subset.drop_duplicates("Chapter")
                .set_index("Chapter")["Product Description"].str[:55].to_dict())
    grp["Product Description"] = grp["Chapter"].map(desc_map).fillna("")

    pivot = grp.pivot_table(
        index=["Chapter", "Sector Bucket", "Product Description"],
        columns="Year", values="Trade Value USD Mn", aggfunc="sum"
    ).reset_index()
    pivot.columns.name = None

    year_cols = [c for c in pivot.columns if str(c).isnumeric() or isinstance(c, (int, float))]
    pivot["Total (USD Mn)"] = pivot[year_cols].sum(axis=1).round(2)
    pivot = pivot.sort_values("Chapter")
    return pivot


def build_unit_price_pivot(master, source_key, config):
    """
    Pivot one flow's unit price (USD per reported quantity unit) by
    HS Code × Year.

    Unit price is computed as sum(Trade Value USD) / sum(Quantity) within
    each (HS Code, Qty Unit, Year) group — a quantity-weighted average,
    not an average of each row's own unit price. That matches how trade
    economists compute "unit value" and keeps a handful of tiny, oddly-
    priced shipments from skewing the result.

    Grouped by HS Code (not the broader Sector Bucket/Chapter used in
    build_pivot) and kept with its Qty Unit, because different products —
    even within the same chapter — can be reported in different quantity
    units (kg vs number of items vs litres); averaging across mismatched
    units would be meaningless. Rows with no usable quantity (qty<=0/blank,
    see build_master_table) are excluded before aggregating.
    """
    subset = master[master["Source File"].str.contains(source_key, case=False, na=False)].copy()
    subset = subset[subset["Quantity"].notna() & (subset["Quantity"] > 0)]
    if subset.empty:
        return pd.DataFrame()

    grp = (subset.groupby(["HS Code", "Qty Unit", "Year"], as_index=False)
           .agg(**{"Trade Value USD": ("Trade Value USD", "sum"),
                   "Quantity":        ("Quantity", "sum")}))
    grp["Unit Price (USD)"] = (grp["Trade Value USD"] / grp["Quantity"]).round(4)

    desc_map   = (subset.drop_duplicates("HS Code")
                  .set_index("HS Code")["Product Description"].str[:55].to_dict())
    bucket_map = (subset.drop_duplicates("HS Code")
                  .set_index("HS Code")["Sector Bucket"].to_dict())
    grp["Product Description"] = grp["HS Code"].map(desc_map).fillna("")
    grp["Sector Bucket"]       = grp["HS Code"].map(bucket_map).fillna("Other")

    pivot = grp.pivot_table(
        index=["HS Code", "Sector Bucket", "Product Description", "Qty Unit"],
        columns="Year", values="Unit Price (USD)", aggfunc="mean"
    ).reset_index()
    pivot.columns.name = None
    pivot = pivot.sort_values("HS Code")
    return pivot


def build_india_share(master, latest_year):
    """India exports / World exports × 100 by chapter and year."""
    we = master[master["Source File"].str.contains("world_exports", case=False, na=False)].copy()
    ie = master[master["Source File"].str.contains("india_exports", case=False, na=False)].copy()

    if we.empty or ie.empty:
        log.warning("  ⚠  Cannot compute India share: missing world or India exports.")
        return pd.DataFrame()

    for df in [we, ie]:
        df["Chapter"] = df["HS Code"].str[:2]

    bucket_map = config["bucket_map"] if False else None  # closure trick — use master's values

    we_grp = we.groupby(["Chapter", "Year"])["Trade Value USD Mn"].sum().reset_index()
    ie_grp = ie.groupby(["Chapter", "Year"])["Trade Value USD Mn"].sum().reset_index()

    merged = we_grp.merge(ie_grp, on=["Chapter", "Year"],
                           suffixes=("_world", "_india"), how="left")
    merged["India Exports (USD Mn)"]  = merged["Trade Value USD Mn_india"].fillna(0)
    merged["World Exports (USD Mn)"]  = merged["Trade Value USD Mn_world"].fillna(0)
    merged["India Share (%)"]         = (
        merged["India Exports (USD Mn)"] / merged["World Exports (USD Mn)"].replace(0, float("nan"))
        * 100
    ).round(2)

    # add bucket and description from we
    desc_map   = we.drop_duplicates("Chapter").set_index("Chapter")["Product Description"].str[:55].to_dict()
    bucket_col = we.drop_duplicates("Chapter").set_index("Chapter")["Sector Bucket"].to_dict()
    merged["Product Description"] = merged["Chapter"].map(desc_map).fillna("")
    merged["Sector Bucket"]       = merged["Chapter"].map(bucket_col).fillna("Other")

    result = merged[["Chapter", "Product Description", "Sector Bucket", "Year",
                      "World Exports (USD Mn)", "India Exports (USD Mn)", "India Share (%)"]
                    ].sort_values(["Chapter", "Year"])
    return result


def build_summary_tables(master, config):
    """
    Returns a dict of named DataFrames for each summary table A–H.
    All values in USD Mn.
    """
    latest = config["end_year"]
    first  = config["start_year"]

    def latest_slice(source_key):
        df = master[master["Source File"].str.contains(source_key, case=False, na=False)].copy()
        df["Chapter"] = df["HS Code"].str[:2]
        return df[df["Year"] == latest]

    we_l = latest_slice("world_exports")
    wm_l = latest_slice("world_imports")
    ie_l = latest_slice("india_exports")
    im_l = latest_slice("india_imports")

    def agg_chapter(df, val_col="Trade Value USD Mn"):
        if df.empty:
            return pd.DataFrame()
        g = df.groupby(["Chapter", "Sector Bucket"], as_index=False)[val_col].sum().round(2)
        desc = df.drop_duplicates("Chapter").set_index("Chapter")["Product Description"].str[:55]
        g["Product Description"] = g["Chapter"].map(desc).fillna("")
        return g

    we_agg = agg_chapter(we_l)
    wm_agg = agg_chapter(wm_l)
    ie_agg = agg_chapter(ie_l)
    im_agg = agg_chapter(im_l)

    # ── A. Total trade by Sector Bucket ─────────────────────
    buckets = sorted(set(
        list(we_agg["Sector Bucket"].unique()) +
        list(ie_agg["Sector Bucket"].unique()) if not we_agg.empty and not ie_agg.empty else []
    ))
    rows_a = []
    for b in buckets:
        we_v = we_agg[we_agg["Sector Bucket"] == b]["Trade Value USD Mn"].sum() if not we_agg.empty else float("nan")
        wm_v = wm_agg[wm_agg["Sector Bucket"] == b]["Trade Value USD Mn"].sum() if not wm_agg.empty else float("nan")
        ie_v = ie_agg[ie_agg["Sector Bucket"] == b]["Trade Value USD Mn"].sum() if not ie_agg.empty else float("nan")
        im_v = im_agg[im_agg["Sector Bucket"] == b]["Trade Value USD Mn"].sum() if not im_agg.empty else float("nan")
        share = round(ie_v / we_v * 100, 2) if we_v and we_v > 0 else float("nan")
        rows_a.append({"Sector Bucket": b,
                        f"World Exports {latest} (USD Mn)": round(we_v, 2),
                        f"World Imports {latest} (USD Mn)": round(wm_v, 2),
                        f"India Exports {latest} (USD Mn)": round(ie_v, 2),
                        f"India Imports {latest} (USD Mn)": round(im_v, 2),
                        "India Share in World Exports %":   share})
    tbl_a = pd.DataFrame(rows_a)

    # ── B. India exports by HS category ─────────────────────
    def merge_ie_we(ie_a, we_a):
        if ie_a.empty:
            return pd.DataFrame()
        m = ie_a.rename(columns={"Trade Value USD Mn": f"India Exports {latest} (USD Mn)"})
        if not we_a.empty:
            m = m.merge(we_a[["Chapter", "Trade Value USD Mn"]].rename(
                columns={"Trade Value USD Mn": f"World Exports {latest} (USD Mn)"}),
                on="Chapter", how="left")
        else:
            m[f"World Exports {latest} (USD Mn)"] = float("nan")
        m["India Share in World Exports %"] = (
            m[f"India Exports {latest} (USD Mn)"] /
            m[f"World Exports {latest} (USD Mn)"].replace(0, float("nan")) * 100
        ).round(2)
        return m.sort_values(f"India Exports {latest} (USD Mn)", ascending=False).reset_index(drop=True)

    tbl_b_raw = merge_ie_we(ie_agg, we_agg)
    if not tbl_b_raw.empty:
        tbl_b_raw.insert(0, "Rank", range(1, len(tbl_b_raw) + 1))
        tbl_b = tbl_b_raw[["Rank", "Chapter", "Product Description", "Sector Bucket",
                             f"India Exports {latest} (USD Mn)",
                             f"World Exports {latest} (USD Mn)",
                             "India Share in World Exports %"]]
        tbl_b.columns = ["Rank", "HS Chapter", "Product Description", "Sector Bucket",
                          f"India Exports {latest} (USD Mn)",
                          f"World Exports {latest} (USD Mn)",
                          "India Share in World Exports %"]
    else:
        tbl_b = pd.DataFrame()

    # ── C. Global exports by HS category ────────────────────
    if not we_agg.empty:
        tbl_c = we_agg.rename(columns={"Trade Value USD Mn": f"World Exports {latest} (USD Mn)"})
        if not ie_agg.empty:
            tbl_c = tbl_c.merge(ie_agg[["Chapter", "Trade Value USD Mn"]].rename(
                columns={"Trade Value USD Mn": f"India Exports {latest} (USD Mn)"}), on="Chapter", how="left")
        else:
            tbl_c[f"India Exports {latest} (USD Mn)"] = float("nan")
        tbl_c["India Share in World Exports %"] = (
            tbl_c.get(f"India Exports {latest} (USD Mn)", 0) /
            tbl_c[f"World Exports {latest} (USD Mn)"].replace(0, float("nan")) * 100
        ).round(2)
        tbl_c = tbl_c.sort_values(f"World Exports {latest} (USD Mn)", ascending=False).reset_index(drop=True)
        tbl_c.insert(0, "Rank", range(1, len(tbl_c) + 1))
    else:
        tbl_c = pd.DataFrame()

    # ── D. India imports by HS category ─────────────────────
    if not im_agg.empty:
        tbl_d = im_agg.rename(columns={"Trade Value USD Mn": f"India Imports {latest} (USD Mn)"}).copy()
        if not ie_agg.empty:
            tbl_d = tbl_d.merge(ie_agg[["Chapter", "Trade Value USD Mn"]].rename(
                columns={"Trade Value USD Mn": f"India Exports {latest} (USD Mn)"}), on="Chapter", how="left")
        else:
            tbl_d[f"India Exports {latest} (USD Mn)"] = float("nan")
        tbl_d["Net Exports (USD Mn)"] = (
            tbl_d.get(f"India Exports {latest} (USD Mn)", 0).fillna(0) -
            tbl_d[f"India Imports {latest} (USD Mn)"].fillna(0)
        ).round(2)
        tbl_d = tbl_d.sort_values(f"India Imports {latest} (USD Mn)", ascending=False).reset_index(drop=True)
        tbl_d.insert(0, "Rank", range(1, len(tbl_d) + 1))
    else:
        tbl_d = pd.DataFrame()

    # ── E. India share in world exports ─────────────────────
    if not ie_agg.empty and not we_agg.empty:
        tbl_e = ie_agg.merge(we_agg[["Chapter", "Trade Value USD Mn"]].rename(
            columns={"Trade Value USD Mn": f"World Exports {latest} (USD Mn)"}), on="Chapter", how="left")
        tbl_e.rename(columns={"Trade Value USD Mn": f"India Exports {latest} (USD Mn)"}, inplace=True)
        tbl_e["India Share in World Exports %"] = (
            tbl_e[f"India Exports {latest} (USD Mn)"] /
            tbl_e[f"World Exports {latest} (USD Mn)"].replace(0, float("nan")) * 100
        ).round(2)
        tbl_e = tbl_e.sort_values("India Share in World Exports %", ascending=False).reset_index(drop=True)
    else:
        tbl_e = pd.DataFrame()

    # ── F & G. CAGR ─────────────────────────────────────────
    def compute_cagr(master_df, source_key, first_yr, last_yr):
        df = master_df[master_df["Source File"].str.contains(source_key, case=False, na=False)].copy()
        if df.empty:
            return pd.DataFrame()
        df["Chapter"] = df["HS Code"].str[:2]
        fy = df[df["Year"] == first_yr].groupby("Chapter")["Trade Value USD Mn"].sum()
        ly = df[df["Year"] == last_yr].groupby("Chapter")["Trade Value USD Mn"].sum()
        n  = last_yr - first_yr
        # No data for the reference first year (e.g. a chapter/sector with no
        # recorded trade in the earliest year of the selected range), or the
        # range collapses to a single year — CAGR can't be computed. Return an
        # empty table instead of crashing the whole clean step (this used to
        # raise "array length 0 does not match index length 1" because `desc`/
        # `bkt` below were built from a Series index that wasn't guaranteed to
        # match fy.index, including whenever fy was empty).
        if fy.empty or n <= 0:
            return pd.DataFrame()
        cagr = ((ly / fy.replace(0, float("nan"))) ** (1 / n) - 1) * 100
        desc = df.drop_duplicates("Chapter").set_index("Chapter")["Product Description"].str[:55]
        bkt  = df.drop_duplicates("Chapter").set_index("Chapter")["Sector Bucket"]
        result = pd.DataFrame({
            "HS Chapter":     fy.index,
            "Product Description": desc.reindex(fy.index).values,
            "Sector Bucket":  bkt.reindex(fy.index).values,
            "First Year":     first_yr,
            "First Year Value (USD Mn)": fy.values.round(2),
            "Latest Year":    last_yr,
            "Latest Year Value (USD Mn)": ly.reindex(fy.index).values.round(2),
            "CAGR %":         cagr.reindex(fy.index).round(2).values,
        }).dropna(subset=["CAGR %"]).sort_values("CAGR %", ascending=False).reset_index(drop=True)
        return result

    tbl_f = compute_cagr(master, "india_exports", first, latest)
    tbl_g = compute_cagr(master, "world_exports", first, latest)

    # ── H. Import dependence ─────────────────────────────────
    if not im_agg.empty:
        tbl_h = im_agg.rename(columns={"Trade Value USD Mn": f"India Imports {latest} (USD Mn)"}).copy()
        if not ie_agg.empty:
            tbl_h = tbl_h.merge(ie_agg[["Chapter", "Trade Value USD Mn"]].rename(
                columns={"Trade Value USD Mn": f"India Exports {latest} (USD Mn)"}), on="Chapter", how="left")
        else:
            tbl_h[f"India Exports {latest} (USD Mn)"] = float("nan")
        tbl_h["Net Exports (USD Mn)"] = (
            tbl_h.get(f"India Exports {latest} (USD Mn)", 0).fillna(0) -
            tbl_h[f"India Imports {latest} (USD Mn)"].fillna(0)
        ).round(2)
        tbl_h["Possible interpretation"] = tbl_h["Net Exports (USD Mn)"].apply(
            lambda x: "Net exporter — India has competitive position" if x > 0
                       else "Net importer — potential domestic manufacturing opportunity" if x < 0
                       else "Balanced")
        tbl_h = tbl_h.sort_values(f"India Imports {latest} (USD Mn)", ascending=False).reset_index(drop=True)
    else:
        tbl_h = pd.DataFrame()

    return {"A_Basket_Summary": tbl_a, "B_India_Exports": tbl_b, "C_Global_Exports": tbl_c,
            "D_India_Imports": tbl_d, "E_India_Share": tbl_e,
            "F_India_CAGR": tbl_f, "G_Global_CAGR": tbl_g, "H_Import_Dependence": tbl_h}



# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# DGCIS HS6 HARMONISATION MODULE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _clean_itchs_code(code_raw):
    """
    Clean an ITCHS/HS code string and return (cleaned_code, length, issue_flag).
    - Remove spaces, dots, hyphens, commas
    - Preserve leading zeros (treat as text)
    - Return length and any QA issue found
    """
    if pd.isna(code_raw) or str(code_raw).strip() in ("", "nan", "None"):
        return None, 0, "Blank or missing code"
    code = str(code_raw).strip()
    # Remove common formatting characters
    code = re.sub(r'[\s\.\-,]', '', code)
    # Strip trailing decimal zeros (e.g. '57011010.0' → '57011010')
    code = re.sub(r'\.0+$', '', code)
    if not code:
        return None, 0, "Empty after cleaning"
    if not code.isdigit():
        return code, len(code), f"Non-numeric characters: {code}"
    length = len(code)
    if length < 6:
        return code.zfill(length), length, f"Code shorter than 6 digits ({length} digits)"
    return code, length, None


def _make_hs6(code, length):
    """Return the HS6 (first 6 digits) from a cleaned code string."""
    if code is None or length < 6:
        return None
    return code[:6]


def load_dgcis_mapping(mapping_file, hs2_filter=None):
    """
    Load the DGCIS QE/PC/ITCHS mapping file and build a harmonised HS6 DataFrame.

    Parameters
    ----------
    mapping_file  : str  — path to QE_PC_HS_Mapping xlsx
    hs2_filter    : list of str — 2-digit HS chapters to keep (None = keep all)

    Returns
    -------
    DataFrame with columns:
      Original Code, Original Code Length, HS2, HS6,
      ITCHS Description, Major Commodity Group (DGCIS Sector),
      Principle Commodity Group (DGCIS Sub-sector),
      QA Flag, Notes
    """
    if not os.path.exists(mapping_file):
        log.warning(f"  ⚠  DGCIS mapping file not found: {mapping_file}")
        return pd.DataFrame()

    try:
        xl = pd.ExcelFile(mapping_file)
        sheet = "EXPORT" if "EXPORT" in xl.sheet_names else xl.sheet_names[0]
        df = pd.read_excel(xl, sheet, header=1, dtype=str)
        df.columns = [str(c).strip() for c in df.columns]
        df = df.dropna(subset=["ITCHS"])
        log.info(f"  ✓  Loaded {len(df)} rows from DGCIS mapping ({sheet} sheet)")
    except Exception as e:
        log.error(f"  ✗  Could not load DGCIS mapping: {e}")
        return pd.DataFrame()

    # hs2_filter (config["hs_codes"]) holds the codes at whatever granularity the
    # sector config was generated with (--hs_level 2/4/6/8), already zero-padded
    # to that length by create_sector_config_from_qe.py — so the chapter is
    # always its first 2 characters. zfill(2) here is only a defensive pad for
    # the rare case of a 1-digit chapter losing its leading zero; it's a no-op
    # for anything already ≥2 chars, so this is safe at every hs_level.
    # (Previously this did `zfill(6)[:2]`, which works for 6-digit HS6 codes but
    # zero-pads shorter codes — e.g. 2-digit chapter codes from hs_level=2
    # sectors — on the wrong side, turning "84" into "008471"->"00". That broke
    # DGCIS harmonisation for every hs_level=2/4 sector, which is most presets.)
    chapter_filter = {str(c).strip().zfill(2)[:2] for c in hs2_filter} if hs2_filter else None

    rows = []
    for _, row in df.iterrows():
        raw_code = row.get("ITCHS", "")
        desc     = str(row.get("ITCHS Description", "")).strip()
        major    = str(row.get("Major Commodity Groups", "")).strip()
        pc       = str(row.get("Principle Commodity Groups", "")).strip()

        cleaned, length, issue = _clean_itchs_code(raw_code)
        hs6 = _make_hs6(cleaned, length)
        hs2 = cleaned[:2] if cleaned and len(cleaned) >= 2 else None

        # Filter to relevant HS2 chapters if specified
        if chapter_filter and hs2 not in chapter_filter:
            continue

        rows.append({
            "Original Code":               str(raw_code).strip(),
            "Original Code Length":        length,
            "HS2 Chapter":                 hs2,
            "HS6":                         hs6,
            "ITCHS Description":           desc,
            "DGCIS Sector (Major Group)":  major,
            "DGCIS Sub-sector (PC Group)": pc,
            "QA Flag":                     issue or "",
            "Notes":                       "",
        })

    result = pd.DataFrame(rows)
    if result.empty:
        log.warning("  ⚠  DGCIS mapping built: 0 codes matched this sector's HS chapters "
                    "(DGCIS_HS_Mapping / HS6_QA_Checks sheets will be empty/placeholder)")
    else:
        log.info(f"  ✓  DGCIS mapping built: {len(result)} codes ({result['HS6'].nunique()} unique HS6)")
    return result


def build_dgcis_mapping_sheet(master, config, dgcis_df):
    """
    Build the DGCIS_HS_Mapping sheet data: match pipeline HS chapters
    to DGCIS entries, show original ITCHS codes, HS6, sector bucket, etc.
    """
    if dgcis_df is None or dgcis_df.empty:
        return pd.DataFrame([{
            "Note": "DGCIS mapping not available. Place QE_PC_HS_Mapping xlsx in pipeline folder and set dgcis_mapping sheet in config."
        }])

    # config["hs_codes"] holds codes at this sector's configured hs_level
    # (2/4/6/8 digits, already zero-padded) — chapter is always the first 2
    # chars; derive it the same way load_dgcis_mapping() does.
    pipeline_hs2 = {str(c).strip().zfill(2)[:2] for c in config["hs_codes"]}
    bucket_map   = config["bucket_map"]

    # Filter DGCIS to pipeline chapters
    mask = dgcis_df["HS2 Chapter"].isin(pipeline_hs2)
    matched = dgcis_df[mask].copy()

    # Add Sector Bucket from pipeline bucket_map
    matched["Pipeline Sector Bucket"] = matched["HS2 Chapter"].map(bucket_map).fillna("Not in pipeline")

    # Add Include/Review/Exclude from config hs_codes sheet
    all_hs = config.get("all_hs_df", pd.DataFrame())
    incl_col = next((c for c in all_hs.columns if "include" in c.lower()), None)
    if incl_col and not all_hs.empty:
        incl_map = dict(zip(all_hs["HS Code"].str.strip(), all_hs[incl_col].str.strip()))
        matched["Include / Exclude"] = matched["HS2 Chapter"].map(incl_map).fillna("Not specified")
    else:
        matched["Include / Exclude"] = "Not specified"

    result = matched[[
        "Original Code", "Original Code Length", "HS2 Chapter", "HS6",
        "ITCHS Description", "DGCIS Sector (Major Group)", "DGCIS Sub-sector (PC Group)",
        "Pipeline Sector Bucket", "Include / Exclude", "QA Flag", "Notes"
    ]].reset_index(drop=True)

    return result


def build_hs6_qa_sheet(master, config, dgcis_df):
    """
    Build the HS6_QA_Checks sheet: flag code issues, cross-check pipeline vs DGCIS.
    """
    rows = []

    if dgcis_df is None or dgcis_df.empty:
        rows.append({
            "Check":         "DGCIS mapping availability",
            "Category":      "Missing data",
            "HS Code / HS6": "N/A",
            "Issue":         "DGCIS mapping file not loaded — all QA checks skipped",
            "Count":         0,
            "Suggested Action": "Load DGCIS mapping file and re-run clean mode",
        })
        return pd.DataFrame(rows)

    # Same chapter derivation as build_dgcis_mapping_sheet() — config["hs_codes"]
    # can be at any configured hs_level (2/4/6/8 digits); chapter = first 2 chars.
    pipeline_hs2 = {str(c).strip().zfill(2)[:2] for c in config["hs_codes"]}

    # 1. Codes shorter than 6 digits
    short = dgcis_df[dgcis_df["Original Code Length"] < 6]
    if not short.empty:
        for _, r in short.iterrows():
            rows.append({
                "Check": "Code length",
                "Category": "Short code",
                "HS Code / HS6": r["Original Code"],
                "Issue": f"Code has only {r['Original Code Length']} digits (expected 6 or 8)",
                "Count": 1,
                "Suggested Action": "Verify code in official ITCHS list before use",
            })

    # 2. Non-numeric codes
    non_num = dgcis_df[dgcis_df["QA Flag"].str.contains("Non-numeric", na=False)]
    if not non_num.empty:
        rows.append({
            "Check": "Code format",
            "Category": "Non-numeric code",
            "HS Code / HS6": ", ".join(non_num["Original Code"].head(10).tolist()),
            "Issue": f"{len(non_num)} codes with non-numeric characters",
            "Count": len(non_num),
            "Suggested Action": "Verify formatting; remove spaces, dots, hyphens",
        })

    # 3. Missing HS6
    miss_hs6 = dgcis_df[dgcis_df["HS6"].isna() | (dgcis_df["HS6"] == "")]
    if not miss_hs6.empty:
        rows.append({
            "Check": "HS6 completeness",
            "Category": "Missing HS6",
            "HS Code / HS6": ", ".join(miss_hs6["Original Code"].head(10).tolist()),
            "Issue": f"{len(miss_hs6)} codes where HS6 could not be derived",
            "Count": len(miss_hs6),
            "Suggested Action": "Check original codes; codes shorter than 6 digits cannot be mapped to HS6",
        })

    # 4. Duplicate HS6 within pipeline (same HS6 under multiple DGCIS sectors)
    if "DGCIS Sector (Major Group)" in dgcis_df.columns:
        hs6_sectors = (dgcis_df[dgcis_df["HS2 Chapter"].isin(pipeline_hs2)]
                       .groupby("HS6")["DGCIS Sector (Major Group)"].nunique())
        multi = hs6_sectors[hs6_sectors > 1]
        if not multi.empty:
            rows.append({
                "Check": "Sector consistency",
                "Category": "HS6 in multiple DGCIS sectors",
                "HS Code / HS6": ", ".join(multi.index.tolist()[:10]),
                "Issue": f"{len(multi)} HS6 codes appear under multiple DGCIS Major Groups",
                "Count": len(multi),
                "Suggested Action": "Review sector bucket assignment; one HS6 may span multiple industries",
            })

    # 5. Pipeline HS chapters not in DGCIS
    dgcis_hs2 = set(dgcis_df["HS2 Chapter"].dropna().unique())
    missing_from_dgcis = pipeline_hs2 - dgcis_hs2
    if missing_from_dgcis:
        rows.append({
            "Check": "Coverage",
            "Category": "Pipeline HS2 not in DGCIS mapping",
            "HS Code / HS6": ", ".join(sorted(missing_from_dgcis)),
            "Issue": f"{len(missing_from_dgcis)} pipeline HS chapters have no matching DGCIS codes",
            "Count": len(missing_from_dgcis),
            "Suggested Action": "Check if these chapters are covered under a different DGCIS group or are non-tradeable",
        })

    # 6. HS6 codes in DGCIS with no Comtrade data
    if not master.empty:
        comtrade_hs2 = set(master["HS Code"].str[:2].unique())
        dgcis_in_pipeline = dgcis_df[dgcis_df["HS2 Chapter"].isin(pipeline_hs2)]
        dgcis_hs6 = set(dgcis_in_pipeline["HS6"].dropna().unique())
        # We can only check at HS2 level since pipeline runs at chapter level
        dgcis_hs2_in_pipeline = set(dgcis_in_pipeline["HS2 Chapter"].dropna().unique())
        no_comtrade = dgcis_hs2_in_pipeline - comtrade_hs2
        if no_comtrade:
            rows.append({
                "Check": "Coverage",
                "Category": "DGCIS HS2 chapters with no Comtrade data",
                "HS Code / HS6": ", ".join(sorted(no_comtrade)),
                "Issue": f"{len(no_comtrade)} DGCIS chapters returned no Comtrade data in this run",
                "Count": len(no_comtrade),
                "Suggested Action": "Check API response for these chapters; may be excluded from Comtrade coverage",
            })

    if not rows:
        rows.append({
            "Check": "All QA checks",
            "Category": "Passed",
            "HS Code / HS6": "N/A",
            "Issue": "No issues found",
            "Count": 0,
            "Suggested Action": "",
        })

    return pd.DataFrame(rows)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# EXCEL WRITING HELPERS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
_HFILL = PatternFill("solid", fgColor="1B3A5C")
_SFILL = PatternFill("solid", fgColor="2874A6")
_AFILL = PatternFill("solid", fgColor="D6EAF8")
_HFONT = Font(bold=True, color="FFFFFF", size=10, name="Calibri")
_TFONT = Font(bold=True, size=13, color="1B3A5C", name="Calibri")
_RFONT = Font(size=10, name="Calibri")
_BFONT = Font(bold=True, size=10, name="Calibri")
_THIN  = Side(border_style="thin", color="CCCCCC")
_BRD   = Border(left=_THIN, right=_THIN, top=_THIN, bottom=_THIN)


def _ws_title(ws, row, text, ncols):
    ws.cell(row=row, column=1, value=text).font = _TFONT
    if ncols > 1:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=ncols)
    ws.row_dimensions[row].height = 22


def _ws_header(ws, row, cols, fill=None):
    fill = fill or _HFILL
    for i, (lbl, w) in enumerate(cols, 1):
        c = ws.cell(row=row, column=i, value=lbl)
        c.fill = fill; c.font = _HFONT; c.border = _BRD
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        ws.column_dimensions[get_column_letter(i)].width = w
    ws.row_dimensions[row].height = 28


def _fast_write_df(ws, df, start_row=1, title_text=None):
    """Write a DataFrame quickly (header + data, minimal per-cell formatting)."""
    r = start_row
    if title_text:
        _ws_title(ws, r, title_text, len(df.columns))
        r += 1
    for ci, col in enumerate(df.columns, 1):
        c = ws.cell(row=r, column=ci, value=str(col))
        c.fill = _HFILL; c.font = _HFONT; c.border = _BRD
        c.alignment = Alignment(horizontal="center", vertical="center")
        ws.column_dimensions[get_column_letter(ci)].width = min(max(len(str(col)) + 2, 12), 40)
    r += 1
    for i, data_row in enumerate(df.itertuples(index=False)):
        alt = i % 2 == 0
        for ci, val in enumerate(data_row, 1):
            c = ws.cell(row=r, column=ci, value=val)
            c.font = _RFONT; c.border = _BRD
            c.alignment = Alignment(vertical="center")
            if alt:
                c.fill = _AFILL
        r += 1
    return r


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CLEAN MODE — BUILD WORKBOOK
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def run_clean(config):
    log.info("\n" + "█" * 60)
    log.info("  CLEAN MODE — building cleaned workbook")
    log.info("█" * 60)

    frames = load_output_files(config)
    if not frames:
        log.error("  ✗  No output files found. Run --mode full first.")
        sys.exit(1)

    master  = build_master_table(frames, config)
    latest  = config["end_year"]
    log.info(f"  Master table: {master.shape[0]} rows × {master.shape[1]} columns")

    pivot_map = {
        "world_exports": "World_Exports",
        "world_imports": "World_Imports",
        "india_exports": "India_Exports",
        "india_imports": "India_Imports",
    }
    share_df   = build_india_share(master, latest)
    summary    = build_summary_tables(master, config)

    # Source log
    source_rows = []
    for label, df in frames.items():
        src = df["_source_file"].iloc[0] if not df.empty else "N/A"
        source_rows.append({
            "Data table":          label,
            "Source file":         src,
            "Source database":     "UN Comtrade",
            "Date created":        datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "Number of rows loaded": len(df),
            "Unit in raw data":    "USD (primaryValue)",
            "Unit in final output":"USD Mn",
            "Transformation":      "Divided by 1,000,000",
            "Errors":              "",
            "Notes":               "",
        })
    source_log_df = pd.DataFrame(source_rows)

    # Error log
    errors_df = pd.DataFrame(_errors_buffer) if _errors_buffer else pd.DataFrame(
        columns=["Timestamp","Mode","Flow","HS Code","Year","Error type","Error message","Action taken"])

    # ── Write workbook ───────────────────────────────────────
    out_path = config["cleaned_file"]
    with pd.ExcelWriter(out_path, engine="openpyxl") as writer:

        # Master_Table — fast write for large sheet
        master.to_excel(writer, sheet_name="Master_Table", index=False, startrow=1)
        ws = writer.sheets["Master_Table"]
        _ws_title(ws, 1, f"{config['sector_name']} — Master Trade Data Table", len(master.columns))
        for ci in range(1, len(master.columns) + 1):
            c = ws.cell(row=2, column=ci)
            c.fill = _HFILL; c.font = _HFONT; c.border = _BRD
            c.alignment = Alignment(horizontal="center", vertical="center")
            ws.column_dimensions[get_column_letter(ci)].width = 18
        log.info("  ✓  Sheet: Master_Table")

        # Pivot sheets
        for source_key, sheet_name in pivot_map.items():
            piv = build_pivot(master, source_key, config)
            ws2 = writer.book.create_sheet(sheet_name)
            if piv.empty:
                ws2.cell(row=1, column=1, value=f"No data available for {sheet_name}")
            else:
                _fast_write_df(ws2, piv, start_row=1,
                               title_text=f"{sheet_name.replace('_',' ')} — Trade Value (USD Mn) by HS Chapter × Year")
            log.info(f"  ✓  Sheet: {sheet_name}")

        # India_Share_World_Exports
        ws_share = writer.book.create_sheet("India_Share_World_Exports")
        if share_df.empty:
            ws_share.cell(row=1, column=1, value="Insufficient data to compute India share.")
        else:
            _fast_write_df(ws_share, share_df, start_row=1,
                           title_text="India Share in World Exports (%) by HS Chapter and Year")
        log.info("  ✓  Sheet: India_Share_World_Exports")

        # Unit_Prices — quantity-weighted USD-per-unit, one stacked table per
        # flow (kept separate from the $-value pivot sheets above since unit
        # price needs its own Qty Unit column and can't be summed across years
        # the way trade value can).
        ws_up = writer.book.create_sheet("Unit_Prices")
        up_row_cursor = 1
        up_section_labels = {
            "world_exports": "World Exports — Unit Price (USD per reported quantity unit)",
            "world_imports": "World Imports — Unit Price (USD per reported quantity unit)",
            "india_exports": "India Exports — Unit Price (USD per reported quantity unit)",
            "india_imports": "India Imports — Unit Price (USD per reported quantity unit)",
        }
        for source_key, up_label in up_section_labels.items():
            up_df = build_unit_price_pivot(master, source_key, config)
            ws_up.cell(row=up_row_cursor, column=1, value=up_label).font = Font(
                bold=True, size=12, color="1B3A5C", name="Calibri")
            ws_up.row_dimensions[up_row_cursor].height = 20
            up_row_cursor += 1
            if up_df.empty:
                ws_up.cell(row=up_row_cursor, column=1,
                           value="Data not available (no quantity reported for this flow's HS codes)").font = Font(
                    italic=True, size=9, color="888888")
                up_row_cursor += 2
            else:
                up_row_cursor = _fast_write_df(ws_up, up_df, start_row=up_row_cursor)
                up_row_cursor += 2
        log.info("  ✓  Sheet: Unit_Prices")

        # Summary_Tables — multiple sub-tables on one sheet
        ws_sum = writer.book.create_sheet("Summary_Tables")
        section_labels = {
            "A_Basket_Summary":   "A.  Total Trade by Sector Bucket",
            "B_India_Exports":    "B.  India Export Categories Ranked",
            "C_Global_Exports":   "C.  Global Export Categories Ranked",
            "D_India_Imports":    "D.  India Import Categories Ranked",
            "E_India_Share":      "E.  India Share in World Exports",
            "F_India_CAGR":       "F.  Fastest-Growing India Export Categories (CAGR)",
            "G_Global_CAGR":      "G.  Fastest-Growing Global Export Categories (CAGR)",
            "H_Import_Dependence":"H.  Import Dependence / Domestic Opportunity",
        }
        row_cursor = 1
        for key, label in section_labels.items():
            df_s = summary.get(key, pd.DataFrame())
            ws_sum.cell(row=row_cursor, column=1, value=label).font = Font(
                bold=True, size=12, color="1B3A5C", name="Calibri")
            ws_sum.row_dimensions[row_cursor].height = 20
            row_cursor += 1
            if df_s is None or df_s.empty:
                ws_sum.cell(row=row_cursor, column=1, value="Data not available").font = Font(
                    italic=True, size=9, color="888888")
                row_cursor += 2
            else:
                row_cursor = _fast_write_df(ws_sum, df_s, start_row=row_cursor)
                row_cursor += 2
        log.info("  ✓  Sheet: Summary_Tables")

        # Source_Log
        ws_src = writer.book.create_sheet("Source_Log")
        _fast_write_df(ws_src, source_log_df, start_row=1, title_text="Source Log")
        log.info("  ✓  Sheet: Source_Log")

        # Errors_Log
        ws_err = writer.book.create_sheet("Errors_Log")
        _fast_write_df(ws_err, errors_df, start_row=1, title_text="Errors Log")
        log.info("  ✓  Sheet: Errors_Log")

        # ── DGCIS_HS_Mapping sheet ─────────────────────────────
        dgcis_file = None
        # Locate DGCIS mapping file: check config, then look in pipeline folder
        dc = config.get("dgcis_config") or {}
        dgcis_file = dc.get("DGCIS mapping file name", "").strip() or                      dc.get("Mapping file", "").strip()
        if not dgcis_file:
            # Auto-detect in same folder as config
            script_dir = os.path.dirname(os.path.abspath(__file__))
            for candidate in os.listdir(script_dir):
                if "qe" in candidate.lower() and candidate.endswith(".xlsx"):
                    dgcis_file = os.path.join(script_dir, candidate)
                    break
        dgcis_df = load_dgcis_mapping(dgcis_file, hs2_filter=config["hs_codes"]) if dgcis_file else pd.DataFrame()

        ws_dgcis = writer.book.create_sheet("DGCIS_HS_Mapping")
        dgcis_sheet_df = build_dgcis_mapping_sheet(master, config, dgcis_df)
        _fast_write_df(ws_dgcis, dgcis_sheet_df, start_row=1,
                       title_text="DGCIS HS Code Mapping (Original ITCHS → HS6)")
        log.info("  ✓  Sheet: DGCIS_HS_Mapping")

        # ── HS6_QA_Checks sheet ────────────────────────────────
        ws_qa = writer.book.create_sheet("HS6_QA_Checks")
        qa_df = build_hs6_qa_sheet(master, config, dgcis_df)
        _fast_write_df(ws_qa, qa_df, start_row=1,
                       title_text="HS6 QA Checks — DGCIS vs Pipeline Cross-check")
        log.info("  ✓  Sheet: HS6_QA_Checks")

    log.info(f"\n  ✅  Cleaned workbook saved → {out_path}")
    sz = os.path.getsize(out_path) / 1024 / 1024
    log.info(f"      File size: {sz:.1f} MB")

    return master, summary, share_df


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE-READY MODE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def _safe_val(v, fmt="{:.2f}"):
    if pd.isna(v):
        return "Data not available"
    try:
        return fmt.format(float(v))
    except Exception:
        return str(v)


def _top_row(df, col, n=1):
    """Return the top n rows sorted by col descending, safely."""
    if df is None or df.empty or col not in df.columns:
        return pd.DataFrame()
    return df.nlargest(n, col)


def build_key_messages(summary, config):
    """Auto-generate key messages from summary data."""
    latest = config["end_year"]
    rows   = []

    # Message 1: Basket summary
    tbl_a = summary.get("A_Basket_Summary", pd.DataFrame())
    if not tbl_a.empty:
        ie_col = f"India Exports {latest} (USD Mn)"
        we_col = f"World Exports {latest} (USD Mn)"
        top_bk = tbl_a.loc[tbl_a[ie_col].idxmax(), "Sector Bucket"] if ie_col in tbl_a else "N/A"
        top_bk_val = _safe_val(tbl_a[ie_col].max() if ie_col in tbl_a else float("nan"), "{:,.0f}")
        total_ie = _safe_val(tbl_a[ie_col].sum() if ie_col in tbl_a else float("nan"), "{:,.0f}")
        rows.append({
            "Source Table":     "A — Basket Summary",
            "Key Finding":      f"India's agri/sector exports are dominated by the '{top_bk}' basket",
            "Supporting Data":  f"'{top_bk}' basket: USD {top_bk_val} Mn ({latest}); Total India sector exports: USD {total_ie} Mn",
            "Why It Matters":   "Basket concentration signals where value-chain upgrading is most needed",
            "Suggested Chart":  "Stacked column or donut chart by basket",
            "Source Note":      f"UN Comtrade, {latest}. All values USD Mn.",
            "Caveat":           "Chapter-level aggregation; does not reflect HS 6-digit subcategory detail.",
        })

    # Message 2: India's top export category
    tbl_b = summary.get("B_India_Exports", pd.DataFrame())
    ie_col = f"India Exports {latest} (USD Mn)"
    we_col = f"World Exports {latest} (USD Mn)"
    top_b = _top_row(tbl_b, ie_col)
    if not top_b.empty:
        ch   = str(top_b.iloc[0].get("HS Chapter", "N/A"))
        desc = str(top_b.iloc[0].get("Product Description", "N/A"))[:40]
        val  = _safe_val(top_b.iloc[0].get(ie_col, float("nan")), "{:,.0f}")
        sh   = _safe_val(top_b.iloc[0].get("India Share in World Exports %", float("nan")))
        rows.append({
            "Source Table":     "B — India Exports",
            "Key Finding":      f"Chapter {ch} ({desc}) is India's largest export category in this sector",
            "Supporting Data":  f"India exports USD {val} Mn ({latest}); India holds {sh}% of global trade",
            "Why It Matters":   "Identifies where India has built export scale and competitiveness",
            "Suggested Chart":  "Horizontal bar chart — top 10 India export chapters",
            "Source Note":      f"UN Comtrade, {latest}. All values USD Mn.",
            "Caveat":           "World export partner = 0 (World aggregate). Chapter-level only.",
        })

    # Message 3: Global market size
    tbl_c = summary.get("C_Global_Exports", pd.DataFrame())
    top_c = _top_row(tbl_c, we_col)
    if not top_c.empty:
        ch   = str(top_c.iloc[0].get("HS Chapter", "N/A"))
        desc = str(top_c.iloc[0].get("Product Description", "N/A"))[:40]
        val  = _safe_val(top_c.iloc[0].get(we_col, float("nan")), "{:,.0f}")
        ind  = _safe_val(top_c.iloc[0].get(f"India Exports {latest} (USD Mn)", float("nan")), "{:,.0f}")
        rows.append({
            "Source Table":     "C — Global Exports",
            "Key Finding":      f"Chapter {ch} ({desc}) is the world's largest export category in this sector",
            "Supporting Data":  f"Global exports: USD {val} Mn ({latest}); India exports: USD {ind} Mn",
            "Why It Matters":   "Indicates where global demand is largest and India's potential headroom",
            "Suggested Chart":  "Horizontal bar chart — global export categories with India overlay",
            "Source Note":      f"UN Comtrade, {latest}. All values USD Mn.",
            "Caveat":           "Global trade includes intra-regional flows.",
        })

    # Message 4: India's highest global share
    tbl_e = summary.get("E_India_Share", pd.DataFrame())
    top_e = _top_row(tbl_e, "India Share in World Exports %")
    if not top_e.empty:
        ch   = str(top_e.iloc[0].get("Chapter", "N/A"))
        desc = str(top_e.iloc[0].get("Product Description", "N/A"))[:40]
        sh   = _safe_val(top_e.iloc[0].get("India Share in World Exports %", float("nan")))
        rows.append({
            "Source Table":     "E — India Share",
            "Key Finding":      f"Chapter {ch} ({desc}) is where India has the highest global export share",
            "Supporting Data":  f"India holds {sh}% of global exports in Chapter {ch} ({latest})",
            "Why It Matters":   "Highlights India's competitive strengths at a chapter level",
            "Suggested Chart":  "Bar chart — India's % share by chapter",
            "Source Note":      f"UN Comtrade, {latest}. Calculated as India exports / World exports × 100.",
            "Caveat":           "Share can be high in small markets. Check absolute values alongside.",
        })

    # Message 5: Import dependence
    tbl_h = summary.get("H_Import_Dependence", pd.DataFrame())
    im_col = f"India Imports {latest} (USD Mn)"
    top_h = _top_row(tbl_h, im_col)
    if not top_h.empty:
        ch   = str(top_h.iloc[0].get("Chapter", "N/A"))
        desc = str(top_h.iloc[0].get("Product Description", "N/A"))[:40]
        val  = _safe_val(top_h.iloc[0].get(im_col, float("nan")), "{:,.0f}")
        net  = _safe_val(top_h.iloc[0].get("Net Exports (USD Mn)", float("nan")), "{:,.0f}")
        rows.append({
            "Source Table":     "H — Import Dependence",
            "Key Finding":      f"Chapter {ch} ({desc}) represents India's highest import exposure in this sector",
            "Supporting Data":  f"India imports USD {val} Mn ({latest}); Net exports: USD {net} Mn",
            "Why It Matters":   "Signals where domestic manufacturing could reduce import dependence",
            "Suggested Chart":  "Grouped bar (imports vs exports) by chapter",
            "Source Note":      f"UN Comtrade, {latest}. All values USD Mn.",
            "Caveat":           "Net exports = India exports minus India imports. Negative = net importer.",
        })

    # Message 6: Fastest growing
    tbl_f = summary.get("F_India_CAGR", pd.DataFrame())
    top_f = _top_row(tbl_f, "CAGR %")
    if not top_f.empty:
        ch   = str(top_f.iloc[0].get("HS Chapter", "N/A"))
        desc = str(top_f.iloc[0].get("Product Description", "N/A"))[:40]
        cagr = _safe_val(top_f.iloc[0].get("CAGR %", float("nan")))
        rows.append({
            "Source Table":     "F — India CAGR",
            "Key Finding":      f"Chapter {ch} ({desc}) is India's fastest-growing export category",
            "Supporting Data":  f"CAGR: {cagr}% per year ({config['start_year']}–{config['end_year']})",
            "Why It Matters":   "Identifies emerging export strengths and momentum categories",
            "Suggested Chart":  "Bar chart — CAGR by chapter",
            "Source Note":      f"UN Comtrade, {config['start_year']}–{config['end_year']}. CAGR = compound annual growth rate.",
            "Caveat":           "CAGR may be distorted if base year value was unusually low or high.",
        })

    return pd.DataFrame(rows)


def build_slide_plan(summary, config):
    """Generate a 12-slide deck plan."""
    latest = config["end_year"]
    sector = config["sector_name"]
    slug   = config["slug"]

    # Extract a few numbers for use in messages
    tbl_b = summary.get("B_India_Exports", pd.DataFrame())
    tbl_c = summary.get("C_Global_Exports", pd.DataFrame())
    tbl_e = summary.get("E_India_Share",    pd.DataFrame())
    ie_col = f"India Exports {latest} (USD Mn)"
    we_col = f"World Exports {latest} (USD Mn)"

    top_india_ch  = tbl_b.iloc[0]["HS Chapter"]  if not tbl_b.empty and "HS Chapter" in tbl_b else "N/A"
    top_global_ch = tbl_c.iloc[0]["HS Chapter"]  if not tbl_c.empty and "HS Chapter" in tbl_c else "N/A"
    top_share_ch  = tbl_e.iloc[0]["Chapter"]     if not tbl_e.empty else "N/A"
    total_india   = _safe_val(tbl_b[ie_col].sum() if not tbl_b.empty and ie_col in tbl_b else float("nan"), "USD {:,.0f} Mn")

    slides = [
        (1,  f"{sector} Trade Landscape",
             f"Overview of {sector} global and India trade, {config['start_year']}–{config['end_year']}",
             "Sector name, date range, country focus, HS chapters covered",
             "Title slide with key stats: total India exports, global market size",
             f"UN Comtrade, {config['start_year']}–{latest}",
             "Chapter-level aggregation only"),
        (2,  "Sector Scope & HS Code Universe",
             f"This analysis covers {len(config['hs_codes'])} HS chapters across {len(config['bucket_df'])} sector buckets",
             "HS code list from config, sector bucket mapping",
             "Table of HS chapters with bucket classification",
             f"UN Comtrade HS classification; sector config file",
             "Chapter-level only; sub-category detail not shown"),
        (3,  "Global Trade Landscape",
             f"World {sector} exports totalled [X] USD Mn in {latest}",
             f"C — Global Exports table, {latest} values",
             "Horizontal bar chart — world exports by HS chapter",
             f"UN Comtrade, {latest}. USD Mn.",
             "Includes intra-regional trade; reporter coverage varies"),
        (4,  "Largest Global Export Categories",
             f"Chapter {top_global_ch} is the largest global export category in {sector}",
             f"C — top 5 chapters by world export value, {latest}",
             "Top-5 horizontal bar chart with values",
             f"UN Comtrade, {latest}. USD Mn.",
             "World aggregate; individual country shares vary"),
        (5,  "Largest Global Import Demand Categories",
             "Identifies chapters with the highest global import demand",
             f"World_Imports summary, {latest} values",
             "Horizontal bar chart — world imports by chapter",
             f"UN Comtrade, {latest}. USD Mn.",
             "World imports reflect demand; may include re-exports"),
        (6,  "India's Export Position",
             f"India's total {sector} exports: {total_india} ({latest})",
             f"B — India Exports table, {latest}",
             "Horizontal bar chart — India exports by chapter",
             f"UN Comtrade, {latest}. USD Mn.",
             "Chapter-level; India = reporter code 699"),
        (7,  "India's Import Profile",
             "Shows where India is a significant importer in this sector",
             f"D — India Imports table, {latest}",
             "Horizontal bar chart — India imports by chapter",
             f"UN Comtrade, {latest}. USD Mn.",
             "Partner = World (code 0)"),
        (8,  "India's Global Export Share by Chapter",
             f"Chapter {top_share_ch} is where India holds the highest global export share",
             f"E — India Share table, {latest}",
             "Bar chart — India's % share of world exports by chapter",
             f"UN Comtrade, {latest}. India exports / World exports × 100.",
             "Share % can be high in small-volume categories"),
        (9,  "Where India is Strong vs Weak",
             "Highlights chapters with high India share vs large global markets where India is absent",
             "E — India Share; C — Global Exports",
             "2×2 matrix or table: India share vs global market size",
             f"UN Comtrade, {latest}. USD Mn.",
             "Qualitative interpretation; use as starting point for deeper analysis"),
        (10, "Import Dependence & Domestic Opportunity",
             "Shows chapters where India imports significantly more than it exports",
             "H — Import Dependence table",
             "Grouped bar chart (India exports vs imports) by chapter",
             f"UN Comtrade, {latest}. Net exports = India exports minus imports.",
             "Net trade position does not capture re-exports or processing trade"),
        (11, "Export Growth Momentum (CAGR)",
             f"Fastest-growing India export chapters by CAGR, {config['start_year']}–{latest}",
             "F — India CAGR table",
             "Bar chart — CAGR % by chapter, sorted descending",
             f"UN Comtrade, {config['start_year']}–{latest}. CAGR = compound annual growth rate.",
             "CAGR can be volatile with small base values"),
        (12, "Initial Implications",
             "Key data-driven observations to guide report and deck development",
             "Key_Messages sheet — top 5–6 findings",
             "Numbered takeaway cards with supporting data points",
             f"UN Comtrade, {config['start_year']}–{latest}.",
             "All observations are data-driven; sector context should be added by analyst"),
    ]

    cols = ["Slide Number", "Slide Title", "Key Message", "Data Points Used",
            "Suggested Visual", "Source Note", "Caveat / Limitation"]
    return pd.DataFrame(slides, columns=cols)


def build_chart_data(master, summary, config):
    """Build chart-ready tables (one per key chart)."""
    latest = config["end_year"]
    tables = {}

    for key, name in [
        ("world_exports", "Global_Exports_by_Chapter"),
        ("world_imports", "Global_Imports_by_Chapter"),
        ("india_exports", "India_Exports_by_Chapter"),
        ("india_imports", "India_Imports_by_Chapter"),
    ]:
        piv = build_pivot(master, key, config)
        if not piv.empty and latest in piv.columns:
            t = piv[["Chapter", "Product Description", "Sector Bucket", latest]].copy()
            t.columns = ["HS Chapter", "Product Description", "Sector Bucket", f"Trade Value USD Mn ({latest})"]
            tables[name] = t.sort_values(f"Trade Value USD Mn ({latest})", ascending=False)

    tables["India_Share_by_Chapter"] = summary.get("E_India_Share", pd.DataFrame())
    tables["Basket_Summary"]         = summary.get("A_Basket_Summary", pd.DataFrame())
    tables["India_CAGR"]             = summary.get("F_India_CAGR", pd.DataFrame())
    tables["Global_CAGR"]            = summary.get("G_Global_CAGR", pd.DataFrame())
    tables["Import_Dependence"]      = summary.get("H_Import_Dependence", pd.DataFrame())
    return tables


def run_slide_ready(config, master=None, summary=None, share_df=None):
    """Build the slide-ready tables workbook."""
    log.info("\n" + "█" * 60)
    log.info("  SLIDE-READY MODE — building slide tables workbook")
    log.info("█" * 60)

    # If master not passed, reload from cleaned file
    if master is None or master.empty:
        cleaned = config["cleaned_file"]
        if not os.path.exists(cleaned):
            log.error(f"  ✗  Cleaned file not found: {cleaned}. Run --mode clean first.")
            sys.exit(1)
        frames  = load_output_files(config)
        master  = build_master_table(frames, config)
        summary = build_summary_tables(master, config)

    key_msgs    = build_key_messages(summary, config)
    slide_plan  = build_slide_plan(summary, config)
    chart_data  = build_chart_data(master, summary, config)

    # Source notes
    latest = config["end_year"]
    src_notes = []
    for name in chart_data:
        src_notes.append({
            "Table name":          name,
            "Data source":         "UN Comtrade (getFinalData)",
            "Years used":          f"{config['start_year']}–{latest}",
            "Unit":                "USD Mn",
            "Caveat":              "Chapter (2-digit HS) level only. Includes all reporters/partners.",
            "Suggested slide note": f"Source: UN Comtrade, {latest}. All values in USD Mn.",
        })
    source_notes_df = pd.DataFrame(src_notes)

    out_path = config["slide_ready_file"]
    with pd.ExcelWriter(out_path, engine="openpyxl") as writer:

        _fast_write_df(writer.book.create_sheet("Key_Messages"), key_msgs, 1,
                       f"Key Messages — {config['sector_name']} Trade Data")
        writer.book.remove(writer.book.active)   # remove default blank sheet
        log.info("  ✓  Sheet: Key_Messages")

        _fast_write_df(writer.book.create_sheet("Slide_Plan"), slide_plan, 1,
                       f"Slide Plan — {config['sector_name']} Trade Landscape")
        log.info("  ✓  Sheet: Slide_Plan")

        ws_cd = writer.book.create_sheet("Chart_Data")
        r = 1
        for tbl_name, df in chart_data.items():
            if df is None or df.empty:
                ws_cd.cell(row=r, column=1, value=f"[{tbl_name}]  Data not available").font = Font(
                    italic=True, size=9, color="888888")
                r += 2
            else:
                r = _fast_write_df(ws_cd, df, start_row=r, title_text=tbl_name)
                r += 2
        log.info("  ✓  Sheet: Chart_Data")

        _fast_write_df(writer.book.create_sheet("Source_Notes"), source_notes_df, 1, "Source Notes")
        log.info("  ✓  Sheet: Source_Notes")

    log.info(f"\n  ✅  Slide-ready workbook saved → {out_path}")

    # NOTE: this mode used to also call run_pptx(config, master, summary) here,
    # silently generating a PowerPoint draft. The Streamlit interface has no
    # PPTX/deck option anywhere, so PPTX generation is intentionally never
    # triggered from this codepath — slide_ready now only produces the Excel
    # summary-tables workbook above. run_pptx() is still defined below and can
    # be called directly from the command line if a deck is ever needed.

    return key_msgs, slide_plan


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PPTX MODE (optional — requires python-pptx)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def run_pptx(config, master, summary):
    """Build a 10-slide draft PowerPoint. Requires: pip install python-pptx"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError:
        log.warning("  ⚠  python-pptx not installed. Skipping PPTX generation.")
        log.warning("     Install with: pip install python-pptx")
        return

    log.info("\n  Building PPTX draft...")
    latest = config["end_year"]
    sector = config["sector_name"]
    slug   = config["slug"]

    DARK   = RGBColor(0x1B, 0x3A, 0x5C)
    MID    = RGBColor(0x28, 0x74, 0xA6)
    ACCENT = RGBColor(0x1A, 0xBC, 0x9C)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_text(slide, text, left, top, width, height,
                 size=18, bold=False, color=None, align="left", wrap=True):
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color or DARK
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        return txb

    def add_rect(slide, left, top, width, height, fill_color):
        from pptx.util import Emu
        shp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
        shp.line.fill.background()
        return shp

    def source_note(slide, text):
        add_text(slide, f"Source: {text}", 0.3, 7.1, 12.7, 0.35, size=8, color=MID)

    # ── Helper to get top-N chapters ───────────────────────
    def _norm(df):
        """Normalise 'Chapter' → 'HS Chapter' so all tables use the same column name."""
        if not df.empty and "Chapter" in df.columns and "HS Chapter" not in df.columns:
            df = df.rename(columns={"Chapter": "HS Chapter"})
        return df

    tbl_b  = _norm(summary.get("B_India_Exports", pd.DataFrame()))
    tbl_c  = _norm(summary.get("C_Global_Exports", pd.DataFrame()))
    tbl_e  = _norm(summary.get("E_India_Share", pd.DataFrame()))
    tbl_d  = _norm(summary.get("D_India_Imports", pd.DataFrame()))
    tbl_h  = _norm(summary.get("H_Import_Dependence", pd.DataFrame()))
    tbl_f  = _norm(summary.get("F_India_CAGR", pd.DataFrame()))
    ie_col = f"India Exports {latest} (USD Mn)"
    we_col = f"World Exports {latest} (USD Mn)"
    im_col = f"India Imports {latest} (USD Mn)"

    total_ie  = tbl_b[ie_col].sum() if not tbl_b.empty and ie_col in tbl_b else 0
    total_we  = tbl_c[we_col].sum() if not tbl_c.empty and we_col in tbl_c else 0

    # ── SLIDE 1: Title ──────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    add_rect(s1, 0, 0, 13.33, 7.5, DARK)
    add_rect(s1, 0, 0, 0.2, 7.5, ACCENT)
    add_text(s1, f"{sector} Trade Landscape", 0.5, 0.6, 12, 1.2, size=42, bold=True, color=WHITE)
    add_text(s1, f"A Decade in Review  |  {config['start_year']} – {latest}", 0.5, 1.85, 12, 0.6, size=20, color=WHITE)
    # Stat cards
    stats = [
        (f"USD {total_ie:,.0f} Mn", "India Sector Exports"),
        (f"USD {total_we:,.0f} Mn", "Global Sector Exports"),
        (f"{len(config['hs_codes'])} chapters", "HS Scope"),
    ]
    for i, (val, lbl) in enumerate(stats):
        bx = 0.5 + i * 4.2
        add_rect(s1, bx, 2.65, 3.8, 1.8, MID)
        add_text(s1, val, bx+0.1, 2.75, 3.6, 0.9, size=22, bold=True, color=ACCENT, align="center")
        add_text(s1, lbl, bx+0.1, 3.65, 3.6, 0.65, size=11, color=WHITE, align="center")
    add_text(s1, f"Country focus: India  |  Source: UN Comtrade  |  {len(config['hs_codes'])} HS chapters",
             0.5, 4.8, 12, 0.4, size=10, color=WHITE)
    source_note(s1, f"UN Comtrade API, {config['start_year']}–{latest}. All values USD Mn.")

    # ── SLIDE 2: Sector Scope ───────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    add_text(s2, "Sector Scope & HS Code Universe", 0.4, 0.2, 12.5, 0.7, size=26, bold=True)
    add_text(s2, f"{len(config['hs_codes'])} HS chapters  |  {len(config['bucket_df'])} sector buckets  |  {config['start_year']}–{latest}",
             0.4, 0.85, 12.5, 0.35, size=13, color=MID)
    # List buckets with codes
    y = 1.4
    for _, row in config["bucket_df"].iterrows():
        add_rect(s2, 0.4, y, 12.5, 0.55, LGRAY)
        add_text(s2, f"  {row['Sector Bucket']}  |  HS Chapters: {row['Included HS codes']}",
                 0.5, y+0.05, 12.2, 0.45, size=13, bold=True)
        y += 0.62
    source_note(s2, "HS 2017 classification. Chapter-level analysis only.")

    # ── SLIDES 3–4: Charts (India + Global exports) ─────────
    def bar_slide(prs, title_text, subtitle, chart_data_obj, source_txt):
        s = prs.slides.add_slide(blank_layout)
        add_text(s, title_text, 0.4, 0.15, 12.5, 0.6, size=24, bold=True)
        add_text(s, subtitle, 0.4, 0.72, 12.5, 0.3, size=11, color=MID)
        if chart_data_obj:
            chart = s.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.3), Inches(1.05), Inches(12.7), Inches(5.75),
                chart_data_obj
            ).chart
            chart.has_legend = False
        source_note(s, source_txt)
        return s

    def make_chart_data(labels, values, series_name="Value"):
        cd = ChartData()
        cd.categories = labels
        cd.add_series(series_name, [round(v, 2) if not pd.isna(v) else 0 for v in values])
        return cd

    # India exports bar
    if not tbl_b.empty and ie_col in tbl_b:
        top10 = tbl_b.head(10)
        lbl = [f"Ch.{r['HS Chapter']} {r['Product Description'][:25]}" for _, r in top10.iterrows()]
        cd  = make_chart_data(lbl, top10[ie_col].tolist(), f"India Exports {latest} (USD Mn)")
        bar_slide(prs, f"India's Top Export Categories, {latest}", "USD Mn  |  Top 10 HS chapters", cd,
                  f"UN Comtrade, {latest}. India (reporter 699), partner World. USD Mn.")

    # Global exports bar
    if not tbl_c.empty and we_col in tbl_c:
        top10 = tbl_c.head(10)
        lbl = [f"Ch.{r['HS Chapter']} {r['Product Description'][:25]}" for _, r in top10.iterrows()]
        cd  = make_chart_data(lbl, top10[we_col].tolist(), f"World Exports {latest} (USD Mn)")
        bar_slide(prs, f"Global Export Categories, {latest}", "USD Mn  |  Top 10 HS chapters", cd,
                  f"UN Comtrade, {latest}. All reporters, partner World. USD Mn.")

    # India imports bar
    if not tbl_d.empty and im_col in tbl_d:
        top10 = tbl_d.head(10)
        lbl = [f"Ch.{r['HS Chapter']} {r['Product Description'][:25]}" for _, r in top10.iterrows()]
        cd  = make_chart_data(lbl, top10[im_col].tolist(), f"India Imports {latest} (USD Mn)")
        bar_slide(prs, f"India's Import Categories, {latest}", "USD Mn  |  Top 10 HS chapters", cd,
                  f"UN Comtrade, {latest}. India (reporter 699), partner World. USD Mn.")

    # India share bar
    if not tbl_e.empty and "India Share (%)" in tbl_e.columns:
        top10 = tbl_e.head(10)
        ch_col = "HS Chapter" if "HS Chapter" in top10.columns else "Chapter"
        lbl = [f"Ch.{r[ch_col]} {r['Product Description'][:25]}" for _, r in top10.iterrows()]
        cd  = make_chart_data(lbl, top10["India Share (%)"].tolist(), "India Share (%)")
        bar_slide(prs, f"India's Share of Global Exports, {latest}", "%  |  India exports / World exports × 100", cd,
                  f"UN Comtrade, {latest}. India exports as % of world exports.")

    # Import dependence (grouped)
    if not tbl_h.empty and im_col in tbl_h and ie_col.replace("Exports", "Exports") in tbl_h.columns:
        top8 = tbl_h.head(8)
        s_dep = prs.slides.add_slide(blank_layout)
        add_text(s_dep, f"Import Dependence & Domestic Opportunity, {latest}", 0.4, 0.15, 12.5, 0.6, size=24, bold=True)
        add_text(s_dep, "India exports vs imports by HS chapter (USD Mn)  |  Negative net = net importer",
                 0.4, 0.72, 12.5, 0.3, size=11, color=MID)
        exp_col = f"India Exports {latest} (USD Mn)"
        if exp_col in top8.columns:
            cd2 = ChartData()
            ch_col2 = "HS Chapter" if "HS Chapter" in top8.columns else "Chapter"
            cd2.categories = [f"Ch.{r[ch_col2]} {r['Product Description'][:20]}" for _, r in top8.iterrows()]
            cd2.add_series("India Exports", [round(v, 2) if not pd.isna(v) else 0 for v in top8[exp_col].tolist()])
            cd2.add_series("India Imports", [round(v, 2) if not pd.isna(v) else 0 for v in top8[im_col].tolist()])
            s_dep.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(0.3), Inches(1.05), Inches(12.7), Inches(5.75), cd2
            )
        source_note(s_dep, f"UN Comtrade, {latest}. Net exports = India exports minus India imports.")

    # CAGR bar
    if not tbl_f.empty and "CAGR %" in tbl_f.columns:
        top10 = tbl_f.dropna(subset=["CAGR %"]).head(10)
        lbl = [f"Ch.{r['HS Chapter']} {r['Product Description'][:25]}" for _, r in top10.iterrows()]
        cd  = make_chart_data(lbl, top10["CAGR %"].tolist(), "CAGR %")
        bar_slide(prs, f"India Export Growth Momentum (CAGR {config['start_year']}–{latest})",
                  f"% per year  |  CAGR = compound annual growth rate, {config['start_year']}–{latest}", cd,
                  f"UN Comtrade, {config['start_year']}–{latest}. CAGR (%) per year.")

    # Key takeaways
    s_kw = prs.slides.add_slide(blank_layout)
    add_rect(s_kw, 0, 0, 13.33, 7.5, DARK)
    add_rect(s_kw, 0, 0, 0.2, 7.5, ACCENT)
    add_text(s_kw, "Key Takeaways", 0.5, 0.2, 12.5, 0.7, size=30, bold=True, color=WHITE)
    add_text(s_kw, "Data-driven observations from UN Comtrade analysis", 0.5, 0.88, 12.5, 0.4, size=13, color=WHITE)
    # Key messages summary (up to 4)
    key_msgs = build_key_messages(summary, config)
    for i, (_, row) in enumerate(key_msgs.head(4).iterrows()):
        col = i % 2
        r   = i // 2
        bx  = 0.5 + col * 6.3
        by  = 1.5 + r * 2.55
        add_rect(s_kw, bx, by, 6.0, 2.3, MID)
        num_shp = s_kw.shapes.add_shape(9, Inches(bx+0.1), Inches(by+0.1), Inches(0.5), Inches(0.5))
        num_shp.fill.solid(); num_shp.fill.fore_color.rgb = ACCENT
        num_shp.line.fill.background()
        add_text(s_kw, str(i + 1), bx+0.12, by+0.12, 0.46, 0.46, size=16, bold=True, color=DARK, align="center")
        add_text(s_kw, str(row.get("Key Finding", ""))[:70],
                 bx + 0.68, by + 0.1, 5.2, 0.5, size=11, bold=True, color=WHITE)
        add_text(s_kw, str(row.get("Supporting Data", ""))[:120],
                 bx + 0.12, by + 0.72, 5.75, 1.45, size=9.5, color=RGBColor(0xD6, 0xEA, 0xF8))
    source_note(s_kw, f"UN Comtrade API, {config['start_year']}–{latest}. All values USD Mn.")

    prs.save(config["pptx_file"])
    log.info(f"  ✅  PPTX saved → {config['pptx_file']}")



# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# COMPETITOR ANALYSIS MODULE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _top_n(n, config_key, competitor_settings):
    """Get top-N setting from competitor_settings."""
    try:
        return int(competitor_settings.get(config_key, str(n)))
    except Exception:
        return n


def build_competitor_exporters(master, config, is_test=False):
    """
    Top exporter countries from world-exports data in Master_Table.
    Uses Reporter column (each country) with Partner = World.
    """
    cs    = config.get("competitor_settings", {})
    top_n_exp = _top_n(10, "Top exporters to retain", cs)
    latest    = config["end_year"]
    first     = config["start_year"]
    ie_col    = f"India Exports {latest} (USD Mn)"

    # World exports: each reporter country exports to World
    we = master[master["Source File"].str.contains("world_exports", case=False, na=False)].copy()
    if we.empty:
        log.warning("  ⚠  No world exports data for competitor analysis.")
        return pd.DataFrame()

    we["Chapter"] = we["HS Code"].str[:2]

    # Limit to latest year for is_test
    if is_test:
        we = we[we["Year"] == latest]

    # Exclude "World" aggregate reporter (if present)
    we = we[~we["Reporter"].str.lower().str.contains("world", na=False)]

    # Group by Reporter, Chapter, Year
    grp = we.groupby(["Reporter", "Chapter", "Year"], as_index=False)["Trade Value USD Mn"].sum()
    grp["Sector Bucket"] = grp["Chapter"].map(config["bucket_map"]).fillna("Other")

    # World total per Chapter per Year (sum of all reporters)
    world_tot = grp.groupby(["Chapter", "Year"])["Trade Value USD Mn"].sum().reset_index()
    world_tot.rename(columns={"Trade Value USD Mn": "World Export Value USD Mn"}, inplace=True)

    grp = grp.merge(world_tot, on=["Chapter", "Year"], how="left")
    grp["Export Market Share %"] = (grp["Trade Value USD Mn"] /
                                     grp["World Export Value USD Mn"].replace(0, float("nan")) * 100).round(2)
    grp.rename(columns={"Trade Value USD Mn": "Export Value USD Mn"}, inplace=True)

    # Rank within each Chapter-Year
    grp["Rank"] = grp.groupby(["Chapter", "Year"])["Export Value USD Mn"].rank(ascending=False, method="min").astype(int)

    # Keep only top_n exporters per chapter per year
    grp = grp[grp["Rank"] <= top_n_exp].copy()

    # Add product description
    desc_map = we.drop_duplicates("Chapter").set_index("Chapter")["Product Description"].str[:55].to_dict()
    grp["Product Description"] = grp["Chapter"].map(desc_map).fillna("")

    # CAGR per reporter per chapter
    fy = grp[grp["Year"] == first].set_index(["Reporter", "Chapter"])["Export Value USD Mn"]
    ly = grp[grp["Year"] == latest].set_index(["Reporter", "Chapter"])["Export Value USD Mn"]
    n  = latest - first
    if n > 0:
        cagr = ((ly / fy.replace(0, float("nan"))) ** (1 / n) - 1) * 100
        cagr_map = cagr.to_dict()
        grp["CAGR %"] = grp.apply(lambda r: round(cagr_map.get((r["Reporter"], r["Chapter"]), float("nan")), 2), axis=1)
    else:
        grp["CAGR %"] = float("nan")

    # Market share change first → latest
    share_first = grp[grp["Year"] == first].set_index(["Reporter", "Chapter"])["Export Market Share %"]
    share_last  = grp[grp["Year"] == latest].set_index(["Reporter", "Chapter"])["Export Market Share %"]
    share_chg = (share_last - share_first).to_dict()
    grp["Market Share Change (pp)"] = grp.apply(
        lambda r: round(share_chg.get((r["Reporter"], r["Chapter"]), float("nan")), 2), axis=1)

    grp["Sector"] = config["sector_name"]
    grp["Notes"]  = ""

    result = grp[[
        "Sector", "Sector Bucket", "Chapter", "Product Description", "Year",
        "Reporter", "Export Value USD Mn", "World Export Value USD Mn",
        "Export Market Share %", "Rank", "CAGR %", "Market Share Change (pp)", "Notes"
    ]].rename(columns={"Chapter": "HS2 Chapter"}).sort_values(
        ["HS2 Chapter", "Year", "Rank"]).reset_index(drop=True)

    log.info(f"  ✓  Competitor exporters: {len(result)} rows, {result['Reporter'].nunique()} unique countries")
    return result


def build_competitor_importers(master, config, is_test=False):
    """
    Top importer countries from world-imports data in Master_Table.
    Uses Reporter = each country, Partner = World (i.e., total imports).
    """
    cs    = config.get("competitor_settings", {})
    top_n_imp = _top_n(10, "Top importers to retain", cs)
    latest    = config["end_year"]
    first     = config["start_year"]

    wm = master[master["Source File"].str.contains("world_imports", case=False, na=False)].copy()
    if wm.empty:
        log.warning("  ⚠  No world imports data for competitor analysis.")
        return pd.DataFrame()

    wm["Chapter"] = wm["HS Code"].str[:2]

    if is_test:
        wm = wm[wm["Year"] == latest]

    wm = wm[~wm["Reporter"].str.lower().str.contains("world", na=False)]

    grp = wm.groupby(["Reporter", "Chapter", "Year"], as_index=False)["Trade Value USD Mn"].sum()
    grp["Sector Bucket"] = grp["Chapter"].map(config["bucket_map"]).fillna("Other")

    world_tot = grp.groupby(["Chapter", "Year"])["Trade Value USD Mn"].sum().reset_index()
    world_tot.rename(columns={"Trade Value USD Mn": "World Import Value USD Mn"}, inplace=True)

    grp = grp.merge(world_tot, on=["Chapter", "Year"], how="left")
    grp["Import Market Share %"] = (grp["Trade Value USD Mn"] /
                                     grp["World Import Value USD Mn"].replace(0, float("nan")) * 100).round(2)
    grp.rename(columns={"Trade Value USD Mn": "Import Value USD Mn"}, inplace=True)

    grp["Rank"] = grp.groupby(["Chapter", "Year"])["Import Value USD Mn"].rank(ascending=False, method="min").astype(int)
    grp = grp[grp["Rank"] <= top_n_imp].copy()

    desc_map = wm.drop_duplicates("Chapter").set_index("Chapter")["Product Description"].str[:55].to_dict()
    grp["Product Description"] = grp["Chapter"].map(desc_map).fillna("")

    fy = grp[grp["Year"] == first].set_index(["Reporter", "Chapter"])["Import Value USD Mn"]
    ly = grp[grp["Year"] == latest].set_index(["Reporter", "Chapter"])["Import Value USD Mn"]
    n  = latest - first
    if n > 0:
        cagr = ((ly / fy.replace(0, float("nan"))) ** (1 / n) - 1) * 100
        cagr_map = cagr.to_dict()
        grp["CAGR %"] = grp.apply(lambda r: round(cagr_map.get((r["Reporter"], r["Chapter"]), float("nan")), 2), axis=1)
    else:
        grp["CAGR %"] = float("nan")

    share_first = grp[grp["Year"] == first].set_index(["Reporter", "Chapter"])["Import Market Share %"]
    share_last  = grp[grp["Year"] == latest].set_index(["Reporter", "Chapter"])["Import Market Share %"]
    share_chg = (share_last - share_first).to_dict()
    grp["Market Share Change (pp)"] = grp.apply(
        lambda r: round(share_chg.get((r["Reporter"], r["Chapter"]), float("nan")), 2), axis=1)

    grp["Sector"] = config["sector_name"]
    grp["Notes"]  = ""

    result = grp[[
        "Sector", "Sector Bucket", "Chapter", "Product Description", "Year",
        "Reporter", "Import Value USD Mn", "World Import Value USD Mn",
        "Import Market Share %", "Rank", "CAGR %", "Market Share Change (pp)", "Notes"
    ]].rename(columns={"Chapter": "HS2 Chapter"}).sort_values(
        ["HS2 Chapter", "Year", "Rank"]).reset_index(drop=True)

    log.info(f"  ✓  Competitor importers: {len(result)} rows, {result['Reporter'].nunique()} unique countries")
    return result


def build_competitor_market_share(comp_exp, comp_imp, config):
    """
    Summary: first-year vs latest-year market share change by country.
    """
    latest = config["end_year"]
    first  = config["start_year"]
    rows   = []

    for flow, df, val_col, share_col in [
        ("Export", comp_exp, "Export Value USD Mn", "Export Market Share %"),
        ("Import", comp_imp, "Import Value USD Mn", "Import Market Share %"),
    ]:
        if df is None or df.empty:
            continue
        fy_df = df[df["Year"] == first]
        ly_df = df[df["Year"] == latest]
        countries = df["Reporter"].unique()
        for country in countries:
            for bucket in df["Sector Bucket"].unique():
                mask_fy = (fy_df["Reporter"] == country) & (fy_df["Sector Bucket"] == bucket)
                mask_ly = (ly_df["Reporter"] == country) & (ly_df["Sector Bucket"] == bucket)
                fy_val   = fy_df[mask_fy][val_col].sum()
                ly_val   = ly_df[mask_ly][val_col].sum()
                fy_share = fy_df[mask_fy][share_col].mean()
                ly_share = ly_df[mask_ly][share_col].mean()
                if ly_val == 0 and fy_val == 0:
                    continue
                n = latest - first
                cagr = (((ly_val / fy_val) ** (1/n) - 1) * 100) if fy_val > 0 and n > 0 else float("nan")
                _rank_val = ly_df[ly_df["Sector Bucket"] == bucket].groupby("Reporter")[val_col].sum().rank(ascending=False).get(country, float("nan")) if not ly_df[ly_df["Sector Bucket"] == bucket].empty else float("nan")
                rank = int(_rank_val) if not pd.isna(_rank_val) else float("nan")
                rows.append({
                    "Country":                           country,
                    "Flow":                              flow,
                    "Sector Bucket":                     bucket,
                    f"First Year ({first}) Market Share %": round(fy_share, 2),
                    f"Latest Year ({latest}) Market Share %": round(ly_share, 2),
                    "Market Share Change (pp)":          round(ly_share - fy_share, 2) if not pd.isna(fy_share) and not pd.isna(ly_share) else float("nan"),
                    f"First Year ({first}) Trade Value USD Mn": round(fy_val, 2),
                    f"Latest Year ({latest}) Trade Value USD Mn": round(ly_val, 2),
                    "CAGR %":                            round(cagr, 2),
                    f"Rank in Latest Year ({latest})":   rank,
                })

    return pd.DataFrame(rows)


def build_competitor_trends(comp_exp, comp_imp, config):
    """
    Time-series trends: top 5 exporters by chapter over all years.
    Returns dict of DataFrames keyed by chapter.
    """
    latest = config["end_year"]
    rows   = []

    for flow, df, val_col, share_col in [
        ("Export", comp_exp, "Export Value USD Mn", "Export Market Share %"),
        ("Import", comp_imp, "Import Value USD Mn", "Import Market Share %"),
    ]:
        if df is None or df.empty:
            continue

        # Top 5 exporters in latest year per chapter
        ly = df[df["Year"] == latest]
        for ch in df["HS2 Chapter"].unique():
            ch_ly = ly[ly["HS2 Chapter"] == ch].nlargest(5, val_col)
            top5  = set(ch_ly["Reporter"].tolist())
            trend = df[(df["HS2 Chapter"] == ch) & (df["Reporter"].isin(top5))]
            for _, r in trend.iterrows():
                rows.append({
                    "Flow":           flow,
                    "HS2 Chapter":    r["HS2 Chapter"],
                    "Sector Bucket":  r["Sector Bucket"],
                    "Country":        r["Reporter"],
                    "Year":           r["Year"],
                    val_col:          r[val_col],
                    share_col:        r[share_col],
                })

    return pd.DataFrame(rows)


def build_competitor_summary(comp_exp, comp_imp, config):
    """
    Key findings from competitor analysis.
    """
    latest  = config["end_year"]
    first   = config["start_year"]
    rows    = []
    ie_col  = "Export Value USD Mn"
    ii_col  = "Import Value USD Mn"

    def top_country(df, val_col, year):
        if df is None or df.empty:
            return ("N/A", 0)
        ly = df[df["Year"] == year].groupby("Reporter")[val_col].sum()
        if ly.empty:
            return ("N/A", 0)
        top = ly.idxmax()
        return (top, round(ly[top], 2))

    top_exp_country, top_exp_val = top_country(comp_exp, ie_col, latest)
    top_imp_country, top_imp_val = top_country(comp_imp, ii_col, latest)

    rows.append({
        "Key Finding":       f"{top_exp_country} is the largest exporter in this sector ({latest})",
        "Supporting Data Point": f"Total exports: USD {top_exp_val:,.0f} Mn ({latest})",
        "Why It Matters":    "Identifies the global market leader and benchmark for India",
        "Suggested Chart":   "Horizontal bar — top 10 exporters by value",
        "Caveat / Limitation": "Country-level data; HS chapter aggregation only",
        "Source Note":       f"UN Comtrade, {latest}. USD Mn.",
    })

    rows.append({
        "Key Finding":       f"{top_imp_country} is the largest importer in this sector ({latest})",
        "Supporting Data Point": f"Total imports: USD {top_imp_val:,.0f} Mn ({latest})",
        "Why It Matters":    "Identifies the largest import market — potential export destination for India",
        "Suggested Chart":   "Horizontal bar — top 10 importers by value",
        "Caveat / Limitation": "Country-level data; HS chapter aggregation only",
        "Source Note":       f"UN Comtrade, {latest}. USD Mn.",
    })

    # India's rank
    if comp_exp is not None and not comp_exp.empty:
        india_exp = comp_exp[comp_exp["Reporter"].str.lower().str.contains("india", na=False)]
        if not india_exp.empty:
            ly_india = india_exp[india_exp["Year"] == latest]
            ly_all   = comp_exp[comp_exp["Year"] == latest].groupby("Reporter")[ie_col].sum()
            ly_india_val = india_exp[india_exp["Year"] == latest][ie_col].sum()
            india_rank = int(ly_all.rank(ascending=False).get("India", float("nan"))) if "India" in ly_all.index else "Not in top N"
            rows.append({
                "Key Finding":       f"India is ranked #{india_rank} among global exporters in this sector ({latest})",
                "Supporting Data Point": f"India exports: USD {ly_india_val:,.0f} Mn ({latest})",
                "Why It Matters":    "Benchmarks India's position relative to global peers",
                "Suggested Chart":   "Ranked bar or table of top exporters with India highlighted",
                "Caveat / Limitation": "Rank based on countries in dataset; may not include all countries",
                "Source Note":       f"UN Comtrade, {latest}. USD Mn.",
            })

    # Market share concentration (top 5 / top 10)
    if comp_exp is not None and not comp_exp.empty:
        ly = comp_exp[comp_exp["Year"] == latest].groupby("Reporter")[ie_col].sum()
        total = ly.sum()
        top5  = ly.nlargest(5).sum()
        top5_share = round(top5 / total * 100, 1) if total > 0 else 0
        rows.append({
            "Key Finding":       f"Top 5 exporters account for {top5_share}% of global exports ({latest})",
            "Supporting Data Point": f"Top 5: {', '.join(ly.nlargest(5).index.tolist()[:5])}",
            "Why It Matters":    "High concentration signals dominant players India must compete with or collaborate with",
            "Suggested Chart":   "Pie or stacked bar showing top 5 vs rest of world",
            "Caveat / Limitation": "Based on countries with available Comtrade data",
            "Source Note":       f"UN Comtrade, {latest}. USD Mn.",
        })

    return pd.DataFrame(rows)


def run_competitor(config, is_test=False):
    """Run competitor analysis and write 5 sheets to cleaned workbook."""
    mode_lbl = "competitor_test" if is_test else "competitor_full"
    log.info("\n" + "█" * 60)
    log.info(f"  COMPETITOR ANALYSIS — {'TEST (latest year only)' if is_test else 'FULL'}")
    log.info("█" * 60)

    cleaned = config["cleaned_file"]
    if not os.path.exists(cleaned):
        log.error(f"  ✗  Cleaned file not found: {cleaned}. Run --mode clean first.")
        sys.exit(1)

    # Load master from cleaned file
    try:
        master_raw = pd.read_excel(cleaned, sheet_name="Master_Table", header=1, dtype={"HS Code": str})
        master_raw.columns = [str(c).strip() for c in master_raw.columns]
        # Ensure Source File column present (may be named differently)
        sf_col = next((c for c in master_raw.columns if "source" in c.lower() and "file" in c.lower()), None)
        if sf_col and sf_col != "Source File":
            master_raw = master_raw.rename(columns={sf_col: "Source File"})
        log.info(f"  ✓  Loaded Master_Table: {len(master_raw)} rows")
    except Exception as e:
        log.error(f"  ✗  Could not load Master_Table from {cleaned}: {e}")
        sys.exit(1)

    comp_exp     = build_competitor_exporters(master_raw, config, is_test=is_test)
    comp_imp     = build_competitor_importers(master_raw, config, is_test=is_test)
    comp_share   = build_competitor_market_share(comp_exp, comp_imp, config)
    comp_trends  = build_competitor_trends(comp_exp, comp_imp, config)
    comp_summary = build_competitor_summary(comp_exp, comp_imp, config)

    # Append sheets to the cleaned workbook
    try:
        wb = openpyxl.load_workbook(cleaned)
        # Remove existing competitor sheets if present
        for sh in ["Competitor_Exporters", "Competitor_Importers",
                   "Competitor_Market_Share", "Competitor_Trends", "Competitor_Summary"]:
            if sh in wb.sheetnames:
                del wb[sh]
        wb.save(cleaned)
        wb.close()

        with pd.ExcelWriter(cleaned, engine="openpyxl", mode="a") as writer:
            for shname, df in [
                ("Competitor_Exporters",    comp_exp),
                ("Competitor_Importers",    comp_imp),
                ("Competitor_Market_Share", comp_share),
                ("Competitor_Trends",       comp_trends),
                ("Competitor_Summary",      comp_summary),
            ]:
                if df is None or df.empty:
                    ws_c = writer.book.create_sheet(shname)
                    ws_c.cell(row=1, column=1, value="Data not available — check Master_Table and re-run.")
                else:
                    ws_c = writer.book.create_sheet(shname)
                    _fast_write_df(ws_c, df, start_row=1, title_text=shname.replace("_", " "))
                log.info(f"  ✓  Sheet: {shname} ({len(df) if df is not None and not df.empty else 0} rows)")
    except Exception as e:
        log.error(f"  ✗  Could not write competitor sheets: {e}")
        traceback.print_exc()
        return

    log.info(f"\n  ✅  Competitor analysis complete → {cleaned}")
    if is_test:
        cs = config.get("competitor_settings", {})
        top_n = _top_n(10, "Top exporters to retain", cs)
        log.info(f"  This was a TEST run (latest year only, top {top_n} per chapter).")
        log.info(f"  Estimated full API calls for competitor_full: none — uses existing Master_Table data.")
        log.info(f"  Run --mode competitor_full to get full time-series data.\n")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# VISUALS MODE — CONTEXT-SETTING GRAPHS
# Replicates the 8-slide structure of the reference deck
# (food_processing_slides.pptx or equivalent)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

_CONTEXT_GRAPH_SPECS = [
    (1,  "Cover — Trade Overview",
         "Title slide with sector name, date range, total India exports, total global exports",
         "BAR/TITLE",
         "B_India_Exports (total), C_Global_Exports (total)",
         "No transformation — sum totals",
         "Cover slide",
         "N/A",
         "USD Mn",
         "Chapter-level aggregation only",
         "Yes", ""),
    (2,  "At a Glance: Key Statistics",
         "4–5 KPI cards: India exports, Global exports, India global share %, Top export chapter",
         "STAT CARDS",
         "B_India_Exports, C_Global_Exports, E_India_Share",
         "Latest year totals only",
         f"At a Glance",
         "N/A",
         "USD Mn / %",
         "Single-year snapshot",
         "Yes", ""),
    (3,  "Top Export Categories",
         "Top 10 HS chapters by India export value in latest year",
         "BAR_CLUSTERED (horizontal)",
         "B_India_Exports — latest year, top 10",
         "Sort by India export value descending",
         "India's Top Export Categories",
         "HS Chapter (y-axis), USD Mn (x-axis)",
         "USD Mn",
         "Chapter-level only",
         "Yes", ""),
    (4,  "India's Share of World Exports",
         "India export share % by HS chapter, top 10 chapters",
         "BAR_CLUSTERED (horizontal)",
         "E_India_Share — latest year, top 10 by share %",
         "Sort by India Share % descending",
         "Where India Leads: Share of World Exports",
         "HS Chapter (y-axis), % (x-axis)",
         "%",
         "Share can be high in small-volume categories",
         "Yes", ""),
    (5,  "Exports by Sector Basket",
         "Column chart showing India exports and global exports by sector bucket",
         "COLUMN_CLUSTERED",
         "A_Basket_Summary — India Exports + World Exports by basket",
         "Group by sector bucket",
         "Trade by Sector Basket",
         "Basket (x-axis), USD Mn (y-axis)",
         "USD Mn",
         "Basket definitions from pipeline config",
         "Yes", ""),
    (6,  "Trade Balance by Basket",
         "Net exports (India exports minus India imports) by sector bucket",
         "COLUMN_CLUSTERED",
         "A_Basket_Summary — India Exports + India Imports by basket",
         "Calculate Net = India Exports - India Imports",
         "India's Trade Balance by Basket",
         "Basket (x-axis), USD Mn (y-axis)",
         "USD Mn",
         "Negative = net importer; positive = net exporter",
         "Yes", ""),
    (7,  "Export Trend 2014–2023",
         "Line chart showing India export trend by basket over all years",
         "LINE_MARKERS",
         "India_Exports pivot (all years) grouped by Sector Bucket",
         "Aggregate by Sector Bucket × Year",
         "India's Export Trend",
         "Year (x-axis), USD Mn (y-axis)",
         "USD Mn",
         "Based on available years in config",
         "Yes", ""),
    (8,  "Key Takeaways",
         "4 auto-generated key messages with supporting data points",
         "TEXT CARDS",
         "Key_Messages sheet",
         "Top 4 key findings",
         "Key Takeaways",
         "N/A",
         "N/A",
         "Data-driven only; analyst context required",
         "Yes", ""),
]

_SPEC_COLS = [
    "Reference Slide Number", "Reference Chart Title", "Chart Type",
    "Required Data", "Data Source Sheet", "Transformation Required",
    "Output Chart Title", "Axis Labels", "Unit",
    "Caveat / Limitation", "Can Be Generated?", "Reason if No"
]


def build_context_graph_specs(config):
    """Return the Context_Graph_Specs DataFrame."""
    rows = []
    for row in _CONTEXT_GRAPH_SPECS:
        d = dict(zip(_SPEC_COLS, row))
        d["Output Chart Title"] = d["Output Chart Title"]
        rows.append(d)
    return pd.DataFrame(rows)


def run_visuals(config, master=None, summary=None, share_df=None, make_deck=False):
    """
    Generate:
      1. [sector_slug]_context_graph_data.xlsx  — chart data workbook (always)
      2. [sector_slug]_context_setting_graphs.pptx — 6-slide deck cloned from
         "Sector's context setting slides.pptx" with sector-specific chart data
         (only when make_deck=True; the Streamlit app does not use this)

    The chart-data computation is fully sector-agnostic: buckets, competitors,
    and importing markets are all derived dynamically from the config and the
    cleaned master table, so this produces correct data for any sector.

    Returns the dict of {sheet_name: DataFrame} chart-data records so callers
    (the Streamlit app) can render interactive charts directly, with no deck.
    """
    import zipfile, io, re
    from pathlib import Path

    log.info("\n" + "█" * 60)
    log.info("  VISUALS MODE — context-setting graph data" + ("" if make_deck else " (data only, no deck)"))
    log.info("█" * 60)

    # ── Locate template (only required if we're building the deck) ───
    template_path = Path(__file__).parent / "Sector's context setting slides.pptx"
    if make_deck and not template_path.exists():
        log.error(f"  ✗  Template not found: {template_path}")
        log.error("     Place 'Sector's context setting slides.pptx' in this folder, or run with make_deck=False.")
        return None

    slug     = config["slug"]
    sector   = config["sector_name"]
    latest   = int(config["end_year"])
    first    = int(config["start_year"])
    vs       = config.get("visual_settings", {})
    out_pptx = vs.get("Output PowerPoint file name", f"{slug}_context_setting_graphs.pptx")
    out_data = vs.get("Output chart data workbook name", f"{slug}_context_graph_data.xlsx")

    # ── Load or reuse master + summary ───────────────────────────
    if master is None or (hasattr(master, 'empty') and master.empty):
        cleaned = config["cleaned_file"]
        if not os.path.exists(cleaned):
            log.error(f"  ✗  Cleaned file not found: {cleaned}. Run --mode clean first.")
            return
        frames  = load_output_files(config)
        master  = build_master_table(frames, config)
        summary = build_summary_tables(master, config)

    tbl_a = summary.get("A_Basket_Summary", pd.DataFrame())

    # ── Helper functions ─────────────────────────────────────────
    def _safe_sum(df):
        if df is None or df.empty or "Trade Value USD Mn" not in df.columns:
            return 0.0
        v = df["Trade Value USD Mn"].sum()
        return 0.0 if pd.isna(v) else float(v)

    def _bn(df):
        """Sum USD Mn and convert to Bn USD (round to 3 dp)."""
        return round(_safe_sum(df) / 1000.0, 3)

    def _ie(yr=None, bkt=None):
        """India exports subset of master."""
        m = master[master["Source File"].str.contains("india_exports", case=False, na=False)]
        if yr  is not None: m = m[m["Year"] == yr]
        if bkt is not None: m = m[m["Sector Bucket"] == bkt]
        return m

    def _we(yr=None, bkt=None, reporter=None):
        """World exports subset of master."""
        m = master[master["Source File"].str.contains("world_exports", case=False, na=False)]
        if yr       is not None: m = m[m["Year"] == yr]
        if bkt      is not None: m = m[m["Sector Bucket"] == bkt]
        if reporter is not None:
            m = m[m["Reporter"].str.contains(reporter[:8], case=False, na=False)]
        return m

    def _wi(yr=None, bkt=None, reporter=None):
        """World imports subset of master (each row = one country's total
        imports from World — Reporter holds the country name, mirroring _we)."""
        m = master[master["Source File"].str.contains("world_imports", case=False, na=False)]
        if yr       is not None: m = m[m["Year"] == yr]
        if bkt      is not None: m = m[m["Sector Bucket"] == bkt]
        if reporter is not None:
            m = m[m["Reporter"].str.contains(reporter[:8], case=False, na=False)]
        return m

    def _cagr(v0, v1, n):
        """Compound annual growth rate (%)."""
        try:
            v0, v1 = float(v0), float(v1)
            if v0 <= 0 or v1 <= 0 or n <= 0:
                return None
            return round(((v1 / v0) ** (1.0 / n) - 1) * 100, 1)
        except Exception:
            return None

    # ── Year anchors ─────────────────────────────────────────────
    years_avail = sorted(master["Year"].dropna().astype(int).unique().tolist())
    yr_end   = latest
    yr_start = first
    mid_t    = (yr_start + yr_end) // 2
    yr_mid   = min(years_avail, key=lambda y: abs(y - mid_t)) if years_avail else yr_start
    yr_2014  = 2014 if 2014 in years_avail else (min(years_avail) if years_avail else yr_start)

    THREE_YEARS = [yr_start, yr_mid, yr_end]
    TWO_YEARS   = [yr_2014, yr_end]
    SIX_YEARS   = sorted(years_avail)[-6:] if len(years_avail) >= 6 else sorted(years_avail)
    while len(SIX_YEARS) < 6:          # pad if fewer than 6 years of data
        SIX_YEARS = [SIX_YEARS[0] - 1] + SIX_YEARS

    log.info(f"  ·  Time anchors — 3-yr: {THREE_YEARS}  2-yr: {TWO_YEARS}  trend: {SIX_YEARS}")

    # ── Sector buckets — top N by latest world exports ────────────
    N_BUCKET_SERIES = 7   # template chart1/chart2 have 7 series
    all_buckets = config.get("sector_buckets", [])
    if not all_buckets:
        all_buckets = sorted(master["Sector Bucket"].dropna().unique().tolist())
    bkt_rank = sorted(all_buckets, key=lambda b: _bn(_we(yr=yr_end, bkt=b)), reverse=True)
    top_buckets = bkt_rank[:N_BUCKET_SERIES]
    while len(top_buckets) < N_BUCKET_SERIES:   # pad if sector has fewer buckets
        top_buckets.append(top_buckets[-1] if top_buckets else "Other")

    # ── Identify top 2 non-China competitors ─────────────────────
    _excl = {"India", "World", "China", "Unspecified", "Areas, nes",
             "Other Asia, nes", "Low-income countries", "High-income countries",
             "Europe", "America", "Africa", "Oceania", "Asia"}

    we_latest = _we(yr=yr_end)
    comp_rank = (
        we_latest[~we_latest["Reporter"].isin(_excl)]
        .groupby("Reporter")["Trade Value USD Mn"].sum()
        .sort_values(ascending=False)
        .index.tolist()
    ) if not we_latest.empty else []

    # Also try competitor sheets if built
    if len(comp_rank) < 2:
        try:
            cdf = pd.read_excel(config["cleaned_file"], sheet_name="Competitor_Exporters")
            if "Year" in cdf.columns and "Reporter" in cdf.columns:
                cdf_l = cdf[cdf["Year"] == yr_end]
                comp_rank = (
                    cdf_l[~cdf_l["Reporter"].isin(_excl)]
                    .groupby("Reporter")["Trade Value USD Mn"].sum()
                    .sort_values(ascending=False)
                    .index.tolist()
                )
        except Exception:
            pass

    comp1 = comp_rank[0] if len(comp_rank) > 0 else "United States of America"
    comp2 = comp_rank[1] if len(comp_rank) > 1 else "Brazil"
    log.info(f"  ·  Top competitors: {comp1} | {comp2}")

    # ── Top 4 importing markets ───────────────────────────────────
    # NOTE: "Partner" in this data source is a flat "World" placeholder for
    # the world_exports/world_imports flows (each row is one Reporter's total
    # vs. World, not a genuine bilateral flow) — grouping by Partner here used
    # to silently match nothing and always fall back to the same hardcoded
    # 4 countries for every sector. The real per-country import totals live
    # in the world_imports flow's Reporter field instead (mirrors how
    # exporters are ranked from world_exports' Reporter field above), so use
    # that — this makes "top importing markets" actual, sector-specific data.
    try:
        wi_latest = _wi(yr=yr_end)
        _excl_p = {"World", "India", "Unspecified", "Other", "Areas, nes", ""}
        imp_rank = (
            wi_latest[~wi_latest["Reporter"].isin(_excl_p)]
            .groupby("Reporter")["Trade Value USD Mn"].sum()
            .sort_values(ascending=False)
        ) if not wi_latest.empty else pd.Series(dtype=float)
        TOP_IMPORTERS = [p for p in imp_rank.index if p and pd.notna(p)][:4]
        if len(TOP_IMPORTERS) < 4:
            raise ValueError("not enough world_imports reporter data")
    except Exception:
        TOP_IMPORTERS = ["European Union", "United States of America", "Japan", "United Kingdom"]
        log.warning("  ⚠  Could not derive top importers from data — using placeholder markets "
                    "(check that the world_imports flow returned data).")

    log.info(f"  ·  Top importers: {TOP_IMPORTERS}")

    # ── Chart 1: Global world exports stacked bar ─────────────────
    #   7 series (buckets) × 3 categories (time points), Bn USD
    chart1_series = []
    for bkt in top_buckets:
        chart1_series.append([_bn(_we(yr=yr, bkt=bkt)) for yr in THREE_YEARS])

    # ── Chart 2: India exports stacked bar ────────────────────────
    #   7 series (buckets) × 2 categories (yr_2014, yr_end), Bn USD
    chart2_series = []
    for bkt in top_buckets:
        chart2_series.append([_bn(_ie(yr=yr, bkt=bkt)) for yr in TWO_YEARS])

    # ── Top bucket for chart3 ─────────────────────────────────────
    top_bucket = top_buckets[0] if top_buckets else ""
    if not tbl_a.empty and "Sector Bucket" in tbl_a.columns:
        ie_col_a = next((c for c in tbl_a.columns if "India Exports" in c), None)
        if ie_col_a:
            tb = tbl_a.dropna(subset=[ie_col_a]).sort_values(ie_col_a, ascending=False)
            if not tb.empty:
                top_bucket = tb.iloc[0]["Sector Bucket"]

    # ── Chart 3: Top bucket — India / comp1 / comp2, 2 time-points ─
    #   2 series (yr_2014, yr_end) × 3 categories (countries), Bn USD
    countries = ["India", comp1, comp2]
    chart3_ser0, chart3_ser1 = [], []
    for ctry in countries:
        if ctry == "India":
            v0, v1 = _bn(_ie(yr=TWO_YEARS[0], bkt=top_bucket)), _bn(_ie(yr=yr_end, bkt=top_bucket))
        else:
            v0, v1 = _bn(_we(yr=TWO_YEARS[0], bkt=top_bucket, reporter=ctry)), _bn(_we(yr=yr_end, bkt=top_bucket, reporter=ctry))
        chart3_ser0.append(v0); chart3_ser1.append(v1)

    # ── Chart 4: Total sector — India / comp1 / comp2, 2 time-points
    #   2 series × 3 categories, Bn USD
    chart4_ser0, chart4_ser1 = [], []
    for ctry in countries:
        if ctry == "India":
            v0, v1 = _bn(_ie(yr=TWO_YEARS[0])), _bn(_ie(yr=yr_end))
        else:
            v0, v1 = _bn(_we(yr=TWO_YEARS[0], reporter=ctry)), _bn(_we(yr=yr_end, reporter=ctry))
        chart4_ser0.append(v0); chart4_ser1.append(v1)

    # ── Chart 5: India/comp1/comp2 shares in top 4 importers (%) ──
    #   3 series × 4 categories
    def _share_in_mkt(ctry, mkt_kw, yr):
        """Export share (%) of ctry in importing market matching mkt_kw."""
        mkt_df = master[(master["Year"] == yr) &
                        (master["Partner"].str.contains(mkt_kw[:8], case=False, na=False))]
        tot = _safe_sum(mkt_df)
        if tot <= 0:
            # Proxy: country global share
            world = _safe_sum(_we(yr=yr))
            if world <= 0:
                return 0.0
            if ctry == "India":
                return round(_safe_sum(_ie(yr=yr)) / world * 100, 1)
            else:
                return round(_safe_sum(_we(yr=yr, reporter=ctry)) / world * 100, 1)
        if ctry == "India":
            num = _safe_sum(mkt_df[mkt_df["Source File"].str.contains("india_exports", case=False, na=False)])
        else:
            num = _safe_sum(mkt_df[mkt_df["Reporter"].str.contains(ctry[:8], case=False, na=False)])
        return round(num / tot * 100, 1) if tot > 0 else 0.0

    chart5_india = [_share_in_mkt("India", imp.split()[0], yr_end) for imp in TOP_IMPORTERS]
    chart5_c1    = [_share_in_mkt(comp1,   imp.split()[0], yr_end) for imp in TOP_IMPORTERS]
    chart5_c2    = [_share_in_mkt(comp2,   imp.split()[0], yr_end) for imp in TOP_IMPORTERS]

    # ── Chart 6: Total imports of top 4 markets (Bn USD) ──────────
    #   1 series × 4 categories
    #   Uses Reporter (not Partner) on the world_imports flow — see the
    #   TOP_IMPORTERS note above for why Partner-based lookups always return
    #   nothing under the default flow configuration.
    chart6_vals = [_bn(_wi(yr=yr_end, reporter=imp)) for imp in TOP_IMPORTERS]

    # ── Chart 7: China shares in top 4 importers (%) ─────────────
    #   1 series × 4 categories
    china_tags = ["China", "156"]
    chart7_vals = []
    for imp in TOP_IMPORTERS:
        mkt_df = master[(master["Year"] == yr_end) &
                        (master["Partner"].str.contains(imp.split()[0][:8], case=False, na=False))]
        tot = _safe_sum(mkt_df)
        china_num = _safe_sum(mkt_df[mkt_df["Reporter"].isin(china_tags)])
        if tot > 0 and china_num > 0:
            chart7_vals.append(round(china_num / tot * 100, 1))
        else:
            # Proxy: China global share
            world = _safe_sum(_we(yr=yr_end))
            china_g = _safe_sum(_we(yr=yr_end, reporter="China"))
            chart7_vals.append(round(china_g / world * 100, 1) if world > 0 else 0.0)

    # ── Chart 8: Market share trends (%) ─────────────────────────
    #   4 series (China/India/comp1/comp2) × 6 years
    china_sh, india_sh, c1_sh, c2_sh = [], [], [], []
    for yr in SIX_YEARS:
        world_t = _safe_sum(_we(yr=yr))
        if world_t <= 0:
            world_t = max(_safe_sum(master[master["Year"] == yr]), 1.0)
        china_t  = _safe_sum(master[(master["Year"] == yr) & master["Reporter"].isin(china_tags)])
        india_t  = _safe_sum(_ie(yr=yr))
        c1_t     = _safe_sum(_we(yr=yr, reporter=comp1))
        c2_t     = _safe_sum(_we(yr=yr, reporter=comp2))
        china_sh.append(round(china_t  / world_t * 100, 1))
        india_sh.append(round(india_t  / world_t * 100, 1))
        c1_sh.append(   round(c1_t     / world_t * 100, 1))
        c2_sh.append(   round(c2_t     / world_t * 100, 1))

    # ── Extended rankings & country trend data (for Streamlit UI controls:
    #    Top-5/Top-10 selectors and user-driven competitor comparison) ──
    # Use the latest year that actually has data — config["end_year"] can be
    # ahead of what was fetched, which would otherwise zero out these rankings.
    yr_end_eff = yr_end if (years_avail and yr_end in years_avail) else (
        max(years_avail) if years_avail else yr_end
    )

    try:
        _excl_agg_ctry = {"World", "Unspecified", "Areas, nes", "Other Asia, nes",
                           "Low-income countries", "High-income countries",
                           "Europe", "America", "Africa", "Oceania", "Asia",
                           "India", ""}
        we_latest_all = _we(yr=yr_end_eff)
        exp_rank_series = (
            we_latest_all[~we_latest_all["Reporter"].isin(_excl_agg_ctry)]
            .groupby("Reporter")["Trade Value USD Mn"].sum()
            .sort_values(ascending=False)
        ) if not we_latest_all.empty else pd.Series(dtype=float)
        india_total = _safe_sum(_ie(yr=yr_end_eff))
        if india_total > 0:
            exp_rank_series = pd.concat(
                [exp_rank_series, pd.Series({"India": india_total})]
            ).sort_values(ascending=False)
        exp_rank_series = exp_rank_series[exp_rank_series > 0].head(20)
        ranking_exporters = pd.DataFrame({
            "Country": exp_rank_series.index,
            f"Exports {yr_end_eff} (Bn USD)": (exp_rank_series.values / 1000.0).round(3),
        })
    except Exception as e:
        log.warning(f"  ⚠  Could not build extended exporter ranking: {e}")
        ranking_exporters = pd.DataFrame(columns=["Country", f"Exports {yr_end_eff} (Bn USD)"])

    try:
        # "Partner" is a flat "World" placeholder in this data source (each row is a
        # Reporter's total vs. World, not genuine bilateral flows) — so the real
        # per-country import ranking comes from Reporter within the world_imports
        # source, mirroring how exporters are ranked from world_exports above.
        wm_latest_all = master[
            master["Source File"].str.contains("world_imports", case=False, na=False) &
            (master["Year"] == yr_end_eff)
        ]
        imp_rank_series = (
            wm_latest_all[~wm_latest_all["Reporter"].isin(_excl_agg_ctry)]
            .groupby("Reporter")["Trade Value USD Mn"].sum()
            .sort_values(ascending=False)
        ) if not wm_latest_all.empty else pd.Series(dtype=float)
        india_imports_total = _safe_sum(
            master[
                master["Source File"].str.contains("india_imports", case=False, na=False) &
                (master["Year"] == yr_end_eff)
            ]
        )
        if india_imports_total > 0:
            imp_rank_series = pd.concat(
                [imp_rank_series, pd.Series({"India": india_imports_total})]
            ).sort_values(ascending=False)
        imp_rank_series = imp_rank_series[imp_rank_series > 0].head(20)
        ranking_importers = pd.DataFrame({
            "Market": imp_rank_series.index,
            f"Imports {yr_end_eff} (Bn USD)": (imp_rank_series.values / 1000.0).round(3),
        })
    except Exception as e:
        log.warning(f"  ⚠  Could not build extended importer ranking: {e}")
        ranking_importers = pd.DataFrame(columns=["Market", f"Imports {yr_end_eff} (Bn USD)"])

    try:
        ranking_buckets = pd.DataFrame({
            "Sector Bucket": bkt_rank,
            f"World Exports {yr_end_eff} (Bn USD)": [_bn(_we(yr=yr_end_eff, bkt=b)) for b in bkt_rank],
        })
    except Exception as e:
        log.warning(f"  ⚠  Could not build extended sector-bucket ranking: {e}")
        ranking_buckets = pd.DataFrame(columns=["Sector Bucket", f"World Exports {yr_end_eff} (Bn USD)"])

    try:
        trend_years = sorted(years_avail) if years_avail else THREE_YEARS
        trend_countries = ranking_exporters["Country"].tolist()[:15] if not ranking_exporters.empty else []
        if "India" not in trend_countries:
            trend_countries = ["India"] + trend_countries
        trend_rows = []
        for ctry in trend_countries:
            row = {"Country": ctry}
            for yr in trend_years:
                row[str(yr)] = _bn(_ie(yr=yr)) if ctry == "India" else _bn(_we(yr=yr, reporter=ctry))
            trend_rows.append(row)
        country_trend = pd.DataFrame(trend_rows)
    except Exception as e:
        log.warning(f"  ⚠  Could not build country trend table: {e}")
        country_trend = pd.DataFrame(columns=["Country"])

    # ── Write chart data workbook ────────────────────────────────
    chart_data_records = {
        "Chart1_GlobalExports_Bn": pd.DataFrame(
            [{"Bucket": b, **{str(y): chart1_series[i][j] for j, y in enumerate(THREE_YEARS)}}
             for i, b in enumerate(top_buckets)]
        ),
        "Chart2_IndiaExports_Bn": pd.DataFrame(
            [{"Bucket": b, **{str(y): chart2_series[i][j] for j, y in enumerate(TWO_YEARS)}}
             for i, b in enumerate(top_buckets)]
        ),
        "Chart3_TopBucketCountries": pd.DataFrame({
            "Country": countries,
            str(TWO_YEARS[0]): chart3_ser0,
            str(yr_end): chart3_ser1,
        }),
        "Chart4_TotalSectorCountries": pd.DataFrame({
            "Country": countries,
            str(TWO_YEARS[0]): chart4_ser0,
            str(yr_end): chart4_ser1,
        }),
        "Chart5_MarketShares_Pct": pd.DataFrame({
            "Importer": TOP_IMPORTERS,
            "India": chart5_india,
            comp1[:20]: chart5_c1,
            comp2[:20]: chart5_c2,
        }),
        "Chart6_ImporterTotals_Bn": pd.DataFrame({
            "Importer": TOP_IMPORTERS,
            "Total Imports (Bn USD)": chart6_vals,
        }),
        "Chart7_ChinaShares_Pct": pd.DataFrame({
            "Importer": TOP_IMPORTERS,
            "China Share (%)": chart7_vals,
        }),
        "Chart8_MarketShareTrends": pd.DataFrame({
            "Year": SIX_YEARS,
            "China": china_sh,
            "India": india_sh,
            comp1[:20]: c1_sh,
            comp2[:20]: c2_sh,
        }),
        "Ranking_Exporters_Bn": ranking_exporters,
        "Ranking_Importers_Bn": ranking_importers,
        "Ranking_SectorBuckets_Bn": ranking_buckets,
        "CountryTrend_AllSector_Bn": country_trend,
    }

    try:
        with pd.ExcelWriter(out_data, engine="openpyxl") as writer:
            for sheet, df in chart_data_records.items():
                df.to_excel(writer, sheet_name=sheet[:31], index=False)
            log.info(f"  ✅  Chart data workbook → {out_data}")
    except Exception as e:
        log.warning(f"  ⚠  Could not write chart data workbook: {e}")

    if not make_deck:
        log.info(f"  ℹ   Skipping PPTX deck (make_deck=False) — chart data only.")
        return chart_data_records

    # ── Clone template and inject chart data ─────────────────────
    chart_data_map = {
        "chart1.xml": chart1_series,
        "chart2.xml": chart2_series,
        "chart3.xml": [chart3_ser0, chart3_ser1],
        "chart4.xml": [chart4_ser0, chart4_ser1],
        "chart5.xml": [chart5_india, chart5_c1, chart5_c2],
        "chart6.xml": [chart6_vals],
        "chart7.xml": [chart7_vals],
        "chart8.xml": [china_sh, india_sh, c1_sh, c2_sh],
    }

    def _rebuild_chart(xml_bytes, new_series_data):
        """
        Replace <c:numCache> block in each <c:ser> with new data.
        Extra template series are zeroed out; new data is never truncated.
        Axis max/min are cleared to allow auto-scaling.
        """
        text = xml_bytes.decode("utf-8")

        # Remove hard-coded axis bounds so the chart auto-scales
        text = re.sub(r'<c:max val="[^"]+"/>\s*', '', text)
        text = re.sub(r'<c:min val="[^"]+"/>\s*', '', text)

        ser_re  = re.compile(r'(<c:ser>)(.*?)(</c:ser>)', re.DOTALL)
        matches = list(ser_re.finditer(text))
        n_new   = len(new_series_data)

        pieces, prev = [], 0
        for i, m in enumerate(matches):
            pieces.append(text[prev:m.start()])
            vals = new_series_data[i] if i < n_new else [0.0] * len(new_series_data[0])
            pts  = ''.join(
                f'<c:pt idx="{j}"><c:v>{float(v)}</c:v></c:pt>'
                for j, v in enumerate(vals)
            )
            new_nc = (
                f'<c:numCache>'
                f'<c:formatCode>General</c:formatCode>'
                f'<c:ptCount val="{len(vals)}"/>'
                f'{pts}'
                f'</c:numCache>'
            )
            updated_body = re.sub(
                r'<c:numCache>.*?</c:numCache>', new_nc, m.group(2), flags=re.DOTALL
            )
            pieces.append(m.group(1) + updated_body + m.group(3))
            prev = m.end()
        pieces.append(text[prev:])
        return ''.join(pieces).encode('utf-8')

    # Read template → patch charts → save
    with open(template_path, "rb") as fh:
        tmpl_bytes = fh.read()

    in_buf  = io.BytesIO(tmpl_bytes)
    out_buf = io.BytesIO()

    with zipfile.ZipFile(in_buf, 'r') as zin:
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data  = zin.read(item.filename)
                fname = item.filename.rsplit('/', 1)[-1]

                if item.filename.startswith('ppt/charts/') and fname in chart_data_map:
                    try:
                        data = _rebuild_chart(data, chart_data_map[fname])
                        log.info(f"  ✓  Chart data injected: {fname}")
                    except Exception as e:
                        log.warning(f"  ⚠  Could not update {fname}: {e}")

                zout.writestr(item, data)

    with open(out_pptx, "wb") as fh:
        fh.write(out_buf.getvalue())

    log.info(f"  ✅  Context-setting PPTX saved → {out_pptx}")
    log.info(f"       Slide 1 : Cover / title")
    log.info(f"       Slide 2 : Global {sector} exports by bucket ({THREE_YEARS[0]}/{THREE_YEARS[1]}/{THREE_YEARS[2]})")
    log.info(f"       Slide 3 : India exports by bucket ({TWO_YEARS[0]}/{TWO_YEARS[1]})")
    log.info(f"       Slide 4 : India vs {comp1[:20]} vs {comp2[:20]} — top bucket + total")
    log.info(f"       Slide 5 : Market shares in {', '.join(i.split()[0] for i in TOP_IMPORTERS)}")
    log.info(f"       Slide 6 : China/India/{comp1[:15]}/{comp2[:15]} market share trends")
    log.info(f"  ℹ   Slide text labels (years, country names, titles) still show T&A reference text.")
    log.info(f"      Update those manually in PowerPoint or run the text-patch step.")

    return chart_data_records

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ALL MODE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def run_all(config):
    log.info("\n" + "█" * 60)
    log.info("  ALL MODE — setup_check → test → full → clean → competitor → visuals")
    log.info("  Note: stops after test to let you verify data before full pull.")
    log.info("█" * 60)

    ok = run_setup_check(config)
    if not ok:
        log.error("  Setup check failed. Fix issues and re-run.")
        sys.exit(1)

    run_fetch(config, is_test=True)
    log.info("\n" + "█" * 60)
    log.info("  STAGE 1 (TEST) COMPLETE.")
    log.info("  ✅  Check the TEST_out_*.xlsx files: do the rows and values look correct?")
    log.info("  To continue to the full run:")
    log.info(f"      python sector_comtrade_pipeline.py --config {config.get('_config_file','')} --mode full")
    log.info("  Then:")
    log.info(f"      python sector_comtrade_pipeline.py --config {config.get('_config_file','')} --mode clean")
    log.info(f"      python sector_comtrade_pipeline.py --config {config.get('_config_file','')} --mode competitor_test")
    log.info(f"      python sector_comtrade_pipeline.py --config {config.get('_config_file','')} --mode competitor_full")
    log.info(f"      python sector_comtrade_pipeline.py --config {config.get('_config_file','')} --mode visuals")
    log.info("█" * 60)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MAIN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def main():
    parser = argparse.ArgumentParser(
        description="Sector Comtrade Pipeline — reusable config-driven trade data automation",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("--config", required=True,
                        help="Path to the sector config Excel file")
    parser.add_argument("--mode",   required=True,
                        choices=["setup_check", "test", "full", "clean", "slide_ready",
                                "competitor_test", "competitor_full", "visuals",
                                "visuals_data", "all"],
                        help="Pipeline mode to run")
    args = parser.parse_args()

    config = read_config(args.config)
    config["_config_file"] = args.config

    log_file = config.get("error_log_file", "pipeline.log")
    setup_logging(log_file)

    log.info(f"\n  Sector Comtrade Pipeline")
    log.info(f"  Sector  : {config['sector_name']}")
    log.info(f"  Config  : {args.config}")
    log.info(f"  Mode    : {args.mode}")
    log.info(f"  HS codes: {len(config['hs_codes'])} (included)")

    if args.mode == "setup_check":
        run_setup_check(config)

    elif args.mode == "test":
        run_fetch(config, is_test=True)
        log.info(
            f"\n  TEST COMPLETE. Check TEST_out_*.xlsx files.\n"
            f"  If data looks correct, run full mode:\n"
            f"      python sector_comtrade_pipeline.py --config {args.config} --mode full\n"
        )

    elif args.mode == "full":
        run_fetch(config, is_test=False)
        log.info(
            f"\n  FULL RUN COMPLETE.\n"
            f"  Run clean mode next:\n"
            f"      python sector_comtrade_pipeline.py --config {args.config} --mode clean\n"
        )

    elif args.mode == "clean":
        master, summary, share_df = run_clean(config)
        log.info(
            f"\n  CLEAN COMPLETE. Open {config['cleaned_file']}.\n"
            f"  Run slide_ready mode next:\n"
            f"      python sector_comtrade_pipeline.py --config {args.config} --mode slide_ready\n"
        )

    elif args.mode == "slide_ready":
        run_slide_ready(config)
        log.info(
            f"\n  SLIDE-READY COMPLETE. Open {config['slide_ready_file']}.\n"
            f"  Next: run competitor_test or visuals mode.\n"
        )

    elif args.mode == "competitor_test":
        run_competitor(config, is_test=True)
        log.info(
            f"\n  COMPETITOR TEST COMPLETE. Check the Competitor_* sheets in {config['cleaned_file']}.\n"
            f"  If data looks correct, run competitor_full:\n"
            f"      python sector_comtrade_pipeline.py --config {args.config} --mode competitor_full\n"
        )

    elif args.mode == "competitor_full":
        run_competitor(config, is_test=False)
        log.info(
            f"\n  COMPETITOR FULL COMPLETE. Check the Competitor_* sheets in {config['cleaned_file']}.\n"
            f"  Next: run visuals mode:\n"
            f"      python sector_comtrade_pipeline.py --config {args.config} --mode visuals\n"
        )

    elif args.mode == "visuals":
        vs = config.get("visual_settings", {})
        out_pptx = vs.get("Output PowerPoint file name", f"{config['slug']}_context_setting_graphs.pptx")
        out_data = vs.get("Output chart data workbook name", f"{config['slug']}_context_graph_data.xlsx")
        run_visuals(config, make_deck=True)
        log.info(
            f"\n  VISUALS COMPLETE.\n"
            f"  ✓  Chart data workbook: {out_data}\n"
            f"  ✓  Context-setting deck: {out_pptx}\n"
        )

    elif args.mode == "visuals_data":
        vs = config.get("visual_settings", {})
        out_data = vs.get("Output chart data workbook name", f"{config['slug']}_context_graph_data.xlsx")
        run_visuals(config, make_deck=False)
        log.info(
            f"\n  VISUALS DATA COMPLETE.\n"
            f"  ✓  Chart data workbook: {out_data}\n"
            f"  ℹ  No PPTX deck generated — view interactive charts in the Streamlit app.\n"
        )

    elif args.mode == "all":
        run_all(config)

    if _errors_buffer:
        log.warning(f"\n  ⚠  {len(_errors_buffer)} error(s) recorded. See {log_file} for details.")
    log.info(f"\n  Full log → {log_file}\n")


if __name__ == "__main__":
    main()
